from dataclasses import dataclass
from pathlib import Path

import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


@dataclass
class ParsedSection:
    kind: str
    order_index: int
    ref: dict
    text: str


def parse_file(path: Path, mime_type: str | None) -> list[ParsedSection]:
    suffix = path.suffix.lower()
    if suffix == ".pdf" or mime_type == "application/pdf":
        return _parse_pdf(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".pptx":
        return _parse_pptx(path)
    if suffix in {".xlsx", ".xlsm"}:
        return _parse_xlsx(path)
    if suffix in {".txt", ".md"} or (mime_type and mime_type.startswith("text/")):
        return _parse_text(path)
    raise ValueError(f"Unsupported file type: {suffix} ({mime_type})")


def _parse_pdf(path: Path) -> list[ParsedSection]:
    sections: list[ParsedSection] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if not text.strip():
                continue
            sections.append(
                ParsedSection(kind="page", order_index=i - 1, ref={"page": i}, text=text)
            )
    return sections


def _parse_docx(path: Path) -> list[ParsedSection]:
    doc = Document(path)
    blocks: list[str] = []
    current: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            if current:
                blocks.append("\n".join(current))
                current = []
            continue
        current.append(text)
    if current:
        blocks.append("\n".join(current))
    return [
        ParsedSection(kind="heading", order_index=i, ref={"paragraph_block": i}, text=t)
        for i, t in enumerate(blocks)
    ]


def _parse_pptx(path: Path) -> list[ParsedSection]:
    prs = Presentation(path)
    sections: list[ParsedSection] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                line = "".join(r.text for r in p.runs).strip()
                if line:
                    parts.append(line)
        text = "\n".join(parts)
        if text.strip():
            sections.append(
                ParsedSection(kind="slide", order_index=i - 1, ref={"slide": i}, text=text)
            )
    return sections


def _parse_xlsx(path: Path) -> list[ParsedSection]:
    wb = load_workbook(path, data_only=True, read_only=True)
    sections: list[ParsedSection] = []
    for s_idx, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append("\t".join(cells))
        text = "\n".join(rows)
        if text.strip():
            sections.append(
                ParsedSection(
                    kind="sheet_range",
                    order_index=s_idx,
                    ref={"sheet": sheet_name},
                    text=text,
                )
            )
    return sections


def _parse_text(path: Path) -> list[ParsedSection]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [ParsedSection(kind="page", order_index=0, ref={"page": 1}, text=text)]
