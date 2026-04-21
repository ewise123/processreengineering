import os
import io
import re
import json
import math
import base64
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List
from datetime import date

from fastapi import FastAPI, File, Form, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, Response
from dotenv import load_dotenv
import anthropic

# Document parsers
import csv
import pdfplumber
import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation

load_dotenv()

app = FastAPI(title="POET API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition"],
)

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".xlsx", ".xls", ".csv", ".png", ".jpg", ".jpeg"}
MAX_FILE_SIZE_MB = 20


# ── Document parsers ──────────────────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> str:
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n\n".join(text_parts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if len(parts) > 1:
            slide_texts.append("\n".join(parts))
    return "\n\n".join(slide_texts)


def extract_text_from_excel(file_bytes: bytes) -> str:
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(c) for c in row if c is not None and str(c).strip())
            if row_text:
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


def extract_text_from_csv(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text))
    return "\n".join(" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader)


def extract_text_from_image(file_bytes: bytes, media_type: str) -> str:
    """Use Claude vision to extract text/content from an image file."""
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured.")
    client = anthropic.Anthropic(api_key=api_key)
    img_b64 = base64.b64encode(file_bytes).decode("utf-8")
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=4000,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "image",
                    "source": {"type": "base64", "media_type": media_type, "data": img_b64},
                },
                {
                    "type": "text",
                    "text": (
                        "Extract all text, labels, data, and meaningful content from this image. "
                        "If it is a process map or diagram, describe the steps, flow, and decisions. "
                        "Return the content in plain text, preserving structure where possible."
                    ),
                },
            ],
        }],
    )
    return msg.content[0].text


def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext in (".docx", ".doc"):
        return extract_text_from_docx(file_bytes)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_bytes)
    elif ext == ".txt":
        return file_bytes.decode("utf-8", errors="ignore")
    elif ext in (".xlsx", ".xls"):
        return extract_text_from_excel(file_bytes)
    elif ext == ".csv":
        return extract_text_from_csv(file_bytes)
    elif ext == ".png":
        return extract_text_from_image(file_bytes, "image/png")
    elif ext in (".jpg", ".jpeg"):
        return extract_text_from_image(file_bytes, "image/jpeg")
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")


# ── Process structure extraction via Claude ───────────────────────────────────

STRUCTURE_PROMPT = """You are a business process analyst. Extract the key process steps from the document and represent them using BPMN 2.0 notation standards.

Return ONLY a valid JSON object in this exact format — no markdown, no explanation:
{
  "process_name": "Short name of the process",
  "steps": [
    {
      "id": "step_1",
      "type": "userTask",
      "name": "Imperative verb + object (max 40 chars)",
      "role": "Actor or department performing this step"
    }
  ],
  "gateways": [
    {
      "id": "gw_1",
      "type": "exclusive",
      "name": "Decision question?",
      "after_step": "step_2",
      "yes_label": "Condition met",
      "no_label": "Condition not met",
      "yes_to": "step_3",
      "no_to": "step_4"
    }
  ]
}

BPMN 2.0 RULES — follow these precisely:

TASK TYPES — choose the most accurate for each step:
  "userTask"         — a human performs the activity (default for most manual steps)
  "serviceTask"      — an automated system or IT service performs the activity with no human intervention
  "manualTask"       — physical or offline work performed by a person without a system (e.g. printing, signing paper)
  "businessRuleTask" — applying a business rule, policy check, or automated decision engine
  "sendTask"         — sending a message, email, or notification to an external party
  "receiveTask"      — waiting to receive a message, document, or trigger from an external party

TASK NAMING — imperative verb + object (the action performed, not a noun phrase):
  CORRECT: "Review application", "Submit claim form", "Approve payment", "Send notification"
  WRONG:   "Application review", "Claim form submission", "Payment approval", "Notification sent"
  Do NOT start the name with the actor/role — the actor is shown in the swimlane header.

GATEWAY TYPES:
  "exclusive"  — exactly one outgoing path is taken based on a condition (XOR — use for most decisions)
  "parallel"   — ALL outgoing paths are taken simultaneously (AND — use when work splits into parallel tracks)
  "inclusive"  — one or more outgoing paths are taken (OR — use when multiple combinations are valid)

GATEWAY NAMING — must be a question that the gateway answers:
  CORRECT: "Application complete?", "Approval granted?", "Risk level acceptable?"
  WRONG:   "Check application", "Decision", "Approval"

GATEWAY CONDITIONS — always use exactly "Yes" and "No" as the yes_label and no_label values.
  Do not use any other text (not "Complete", "Approved", "Accepted" — only "Yes" and "No").
  For parallel gateways, omit yes_label and no_label (all paths are always taken).

GATEWAY ROUTING — critical rules for no_to and yes_to:
- "yes_to" is the step the process continues to on the YES path (the normal forward flow).
  It is ALWAYS the step immediately after the gateway in the sequence — do NOT set it explicitly.
- "no_to" MUST point to a DIFFERENT step than the one immediately after the gateway.
  If Yes continues to step_3, No must go somewhere else (e.g. step_5, or omit to route to End).
  NEVER set "no_to" to the same step ID as the step directly after the gateway.
- If you cannot identify a meaningful No destination, omit "no_to" (the system will route No to the End event).
- A gateway with both Yes and No going to the same step is invalid — use a task instead.

ADDITIONAL RULES:
- The number of steps is defined by the detail level instruction below — follow it precisely
- Step IDs must be unique snake_case strings
- The "role" field is REQUIRED for every step — identify the actor, department, or system. If unspecified, use "Process Team"
- Group related steps under the same role name so swimlanes are meaningful
- Include gateways only where there is a clear decision or split in the process
- If no clear branches exist, return an empty gateways array"""


IDENTIFY_PROCESSES_PROMPT = """Analyze the provided document(s) and identify all distinct business processes described within them.

Return ONLY a valid JSON array — no prose, no markdown fences. Each item must have:
- "name": short process name in Title Case (3–6 words)
- "description": one sentence describing what the process achieves

Rules:
- List only genuine end-to-end processes, not individual tasks or sub-steps
- Typical document contains 1–5 processes
- If the content describes one coherent process, return a single-item array

Example output:
[
  {"name": "Invoice Approval Process", "description": "Manages review and sign-off of vendor invoices before payment is released."},
  {"name": "Vendor Onboarding Process", "description": "Guides new suppliers through registration, compliance checks, and system setup."}
]"""


LEVEL_INSTRUCTIONS = {
    "1": (
        "This is a LEVEL 1 — Process Landscape map. Purpose: executive alignment and scope definition.\n"
        "STEP COUNT: 5–10 major process phases only (e.g. 'Receive Claim → Triage → Adjudicate → Settle → Close').\n"
        "CONTENT RULES:\n"
        "  - Show ONLY the major end-to-end stages — what happens, not how.\n"
        "  - Do NOT include sub-tasks, decision gateways, exception paths, or system interactions.\n"
        "  - Return an EMPTY gateways array — no decision points at this level.\n"
        "  - Use a single broad role (e.g. 'Business', 'Operations') or omit swimlane differentiation.\n"
        "  - Think of this as a value chain: simple linear boxes representing major phases.\n"
        "  - Suitable for: steering committees, strategy sessions, executive briefings."
    ),
    "2": (
        "This is a LEVEL 2 — End-to-End Cross-Functional Process map. Purpose: understand flow across teams.\n"
        "STEP COUNT: 15–25 activities.\n"
        "CONTENT RULES:\n"
        "  - Introduce full BPMN elements: start/end events, activities, and key decision gateways.\n"
        "  - Show major steps AND handoffs between functions — who does what and when.\n"
        "  - Use distinct swimlane roles for each function or department involved (e.g. 'Call Center', 'Adjuster', 'Finance').\n"
        "  - Include key decision points (exclusive gateways) that drive the main flow, but omit low-level task detail.\n"
        "  - Do NOT break activities into individual sub-tasks — keep each step at a functional action level.\n"
        "  - Suitable for: process owners, operations managers, transformation programmes."
    ),
    "3": (
        "This is a LEVEL 3 — Detailed Operational Workflow map. Purpose: diagnose inefficiencies and design improvements.\n"
        "STEP COUNT: 30–50 activities.\n"
        "CONTENT RULES:\n"
        "  - Break each Level 2 activity into its granular component tasks.\n"
        "  - Include ALL of the following where they exist:\n"
        "      * Detailed decision logic (exclusive, parallel, and inclusive gateways)\n"
        "      * Exception paths and error handling\n"
        "      * Rework loops (e.g. 'Return for correction → Re-check')\n"
        "      * System interactions (e.g. 'Log in policy system', 'Update CRM record') — but not field-level detail\n"
        "  - Assign specific roles, departments, or systems to EVERY step.\n"
        "  - Capture real operational complexity — this map should expose bottlenecks and handoff delays.\n"
        "  - Suitable for: process improvement teams, Six Sigma, Lean, operational redesign."
    ),
    "4": (
        "This is a LEVEL 4 — Work Instruction / System-Level Detail map. Purpose: execution, training, and automation design.\n"
        "STEP COUNT: 50–80 activities (use the most important if the process would exceed 80).\n"
        "CONTENT RULES:\n"
        "  - The lowest level of detail — step-by-step instructions for individuals or systems.\n"
        "  - Each step should describe a single atomic action, such as:\n"
        "      * Opening a specific screen or system\n"
        "      * Entering a specific field, selecting a dropdown, or triggering a workflow\n"
        "      * Applying a specific business rule or calculation\n"
        "      * Sending a specific form or notification to a named recipient\n"
        "  - Include all business rules, validation checks, and conditional logic.\n"
        "  - Assign precise roles, named systems, or tools to every step (e.g. 'Claims System', 'RPA Bot', 'Senior Adjuster').\n"
        "  - Suitable for: SOPs, RPA development, system configuration, staff training, audit documentation."
    ),
}


def extract_process_structure(client: anthropic.Anthropic, document_text: str, bpmn_level: str = "2", focus_process: str = "", map_type: str = "") -> dict:
    level_note = LEVEL_INSTRUCTIONS.get(bpmn_level, LEVEL_INSTRUCTIONS["2"])
    system_prompt = STRUCTURE_PROMPT + f"\n\nIMPORTANT — Detail level instruction:\n{level_note}"
    truncated = document_text[:8000] if len(document_text) > 8000 else document_text
    focus_note = (f"\n\nFocus exclusively on the process named: \"{focus_process}\". "
                  f"Ignore all other processes described in the document."
                  if focus_process else "")
    if map_type == "current_state":
        focus_note += ("\n\nMAP TYPE — CURRENT STATE: Document the process EXACTLY AS IT EXISTS TODAY. "
                       "Show the actual workflow including manual steps, handoffs, delays, and inefficiencies as they currently occur. "
                       "Do not optimise or idealise — capture reality.")
    elif map_type == "future_state":
        focus_note += ("\n\nMAP TYPE — FUTURE STATE: Design the process AS IT SHOULD WORK after improvement. "
                       "Show the optimised, streamlined workflow with inefficiencies removed, automation where applicable, "
                       "and best-practice steps applied. This is the target desired state.")
    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        system=system_prompt,
        messages=[{
            "role": "user",
            "content": f"Extract the process structure from this document:{focus_note}\n\n{truncated}"
        }]
    )
    raw = message.content[0].text.strip()
    # Strip any accidental markdown fences
    raw = re.sub(r"^```[a-zA-Z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw).strip()
    # Remove control characters that break JSON parsing
    raw = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", raw)
    # If Claude wrapped JSON in prose, extract the first {...} block
    match = re.search(r"\{[\s\S]*\}", raw)
    if match:
        raw = match.group(0)
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Could not parse process structure from Claude response: {e}")


# ── Programmatic BPMN XML generation with swimlanes ──────────────────────────

def _strip_actor(name: str, role: str) -> str:
    """Remove a leading actor/role prefix from a task label.

    The actor is already shown in the swimlane header so prefixing it in the
    task name is redundant.  Handles exact role match as well as individual
    words within the role (e.g. role='Risk Engineering / Actuarial' strips
    leading 'Risk', 'Engineering', or 'Actuarial').
    """
    if not name or not role:
        return name
    # Try full role string first (e.g. "Broker sends…" when role="Broker")
    for candidate in [role] + role.replace('/', ' ').split():
        candidate = candidate.strip()
        if len(candidate) < 3:
            continue
        if name.lower().startswith(candidate.lower()):
            remainder = name[len(candidate):].lstrip(' \t-–:,')
            if remainder:
                return remainder[0].upper() + remainder[1:]
    return name


def build_bpmn_xml(structure: dict) -> str:
    process_name = structure.get("process_name", "Process")
    steps        = structure.get("steps", [])
    gateways     = structure.get("gateways", [])

    # ── Collect unique roles in document order ────────────────────────────────
    role_order: list[str] = []
    _seen: set[str] = set()
    for step in steps:
        r = (step.get("role") or "Process Team").strip()
        if r not in _seen:
            role_order.append(r)
            _seen.add(r)
    if not role_order:
        role_order = ["Process Team"]

    def lane_of(role: str) -> int:
        r = (role or "Process Team").strip()
        return role_order.index(r) if r in _seen else 0

    # ── Layout constants ──────────────────────────────────────────────────────
    POOL_HDR_W  = 30    # pool name strip width (vertical label on far left)
    LANE_LBL_W  = 120   # width reserved for lane label strip inside each lane
    CONTENT_OFF = POOL_HDR_W + LANE_LBL_W + 60  # first element centre x-offset from P_X
    H_STEP      = 175   # horizontal distance between element centres
    LANE_H      = 150   # height of each swim lane
    P_X         = 50    # participant left x
    P_Y         = 50    # participant top y
    MARGIN_R    = 100   # right padding after last element

    # Canonical BPMN 2.0 task element names — all task variants share the same size
    TASK_TYPES = {"userTask", "serviceTask", "manualTask", "businessRuleTask",
                  "sendTask", "receiveTask", "scriptTask", "task"}
    GATEWAY_TYPES = {"exclusive", "parallel", "inclusive"}
    # Map gateway type strings from Claude → BPMN 2.0 element names
    GATEWAY_ELEMENT = {
        "exclusive": "exclusiveGateway",
        "parallel":  "parallelGateway",
        "inclusive": "inclusiveGateway",
    }

    SIZES: dict[str, tuple[int, int]] = {
        "startEvent": (36, 36),
        "endEvent":   (36, 36),
        "task":       (120, 80),   # 80 px tall: 24 px icon zone + 56 px for text
        "gateway":    (50, 50),
    }

    # ── Build ordered element list ────────────────────────────────────────────
    gateway_map = {gw["after_step"]: gw for gw in gateways}
    first_role  = (steps[0].get("role") or "Process Team").strip() if steps else "Process Team"
    last_role   = (steps[-1].get("role") or "Process Team").strip() if steps else "Process Team"

    elements: list[dict] = []
    elements.append({"id": "Start_1", "type": "startEvent", "name": "Start", "role": first_role})
    for step in steps:
        role = (step.get("role") or "Process Team").strip()
        task_name = _strip_actor(step.get("name", ""), role)
        # Honour the step type from Claude; fall back to userTask if absent/unknown
        raw_type = (step.get("type") or "userTask").strip()
        step_type = raw_type if raw_type in TASK_TYPES else "userTask"
        elements.append({"id": step["id"], "type": step_type, "name": task_name, "role": role})
        if step["id"] in gateway_map:
            gw = gateway_map[step["id"]]
            gw_type = GATEWAY_ELEMENT.get((gw.get("type") or "exclusive").strip(), "exclusiveGateway")
            elements.append({"id": gw["id"], "type": gw_type,
                             "name": gw.get("name", "Decision?"),
                             "yes_label": gw.get("yes_label", "Yes"),
                             "no_label":  gw.get("no_label",  "No"),
                             "role": role})
    elements.append({"id": "End_1", "type": "endEvent", "name": "End", "role": last_role})

    # ── Assign positions (all elements share same centre-Y within their lane) ──
    for col, el in enumerate(elements):
        li   = lane_of(el["role"])
        # All task variants share the task size; gateways share the gateway size
        size_key = (el["type"] if el["type"] in SIZES
                    else "task" if el["type"] in TASK_TYPES
                    else "gateway" if "Gateway" in el["type"]
                    else "task")
        w, h = SIZES[size_key]
        cx   = P_X + CONTENT_OFF + col * H_STEP
        cy   = P_Y + li * LANE_H + LANE_H // 2
        el.update({
            "col": col, "lane_idx": li, "w": w, "h": h,
            "cx": cx, "cy": cy,
            "x": cx - w // 2, "y": cy - h // 2,
        })
        el["right_x"] = cx + w // 2
        el["left_x"]  = cx - w // 2
        el["top_y"]   = cy - h // 2
        el["bot_y"]   = cy + h // 2

    el_by_id = {el["id"]: el for el in elements}

    # Participant bounding box
    last_cx = max(el["cx"] for el in elements)
    P_W     = (last_cx - P_X) + SIZES["task"][0] // 2 + MARGIN_R
    P_H     = len(role_order) * LANE_H

    # ── Build sequence flows ──────────────────────────────────────────────────
    flows: list[dict] = []
    fc = [1]
    # Use .get() + `or` so that None / missing no_to both fall back to "End_1"
    gateway_no = {gw["id"]: (gw.get("no_to") or "") for gw in gateways}

    def add_flow(src_id: str, tgt_id: str, name: str = "") -> None:
        fid = f"Flow_{fc[0]}"; fc[0] += 1
        src = el_by_id.get(src_id)
        tgt = el_by_id.get(tgt_id)
        if not (src and tgt):
            return
        forward   = tgt["col"] > src["col"]
        adjacent  = abs(tgt["col"] - src["col"]) == 1
        same_lane = src["lane_idx"] == tgt["lane_idx"]

        # ── Corridor calculation ─────────────────────────────────────────────
        # For skip/loop routing we must use the corridor of whichever lane is
        # furthest in the routing direction, so the final vertical segment always
        # approaches the target from OUTSIDE its bounding box (never through it).
        #
        # Skip-forward → route via the BOTTOM corridor of the lowest lane
        #   (highest lane_idx) involved.  The arrow descends into the corridor,
        #   travels horizontally, then rises to the target's bottom edge.
        #
        # Backward loop → route via the TOP corridor of the highest lane
        #   (lowest lane_idx) involved.  The arrow rises into the corridor,
        #   travels horizontally, then descends to the target's top edge.
        max_lane = max(src["lane_idx"], tgt["lane_idx"])
        min_lane = min(src["lane_idx"], tgt["lane_idx"])
        bot_corridor = P_Y + (max_lane + 1) * LANE_H - 15   # below all elements in max_lane
        top_corridor = P_Y + min_lane * LANE_H + 15          # above all elements in min_lane

        if forward and adjacent and same_lane:
            # Straight horizontal arrow within same lane
            wps = [(src["right_x"], src["cy"]),
                   (tgt["left_x"],  tgt["cy"])]
        elif forward and adjacent:
            # Adjacent elements in different lanes: L-shape at midpoint x
            mid_x = (src["right_x"] + tgt["left_x"]) // 2
            wps = [(src["right_x"], src["cy"]),
                   (mid_x,          src["cy"]),
                   (mid_x,          tgt["cy"]),
                   (tgt["left_x"],  tgt["cy"])]
        elif forward:
            # Skip-forward: descend to bottom corridor of the lowest lane,
            # travel right, then rise to the target's bottom edge — never
            # enters the target task from inside
            wps = [(src["cx"], src["bot_y"]),
                   (src["cx"], bot_corridor),
                   (tgt["cx"], bot_corridor),
                   (tgt["cx"], tgt["bot_y"])]
        else:
            # Backward loop: rise to top corridor of the highest lane,
            # travel left, then descend to the target's top edge — never
            # enters the target task from inside
            wps = [(src["cx"], src["top_y"]),
                   (src["cx"], top_corridor),
                   (tgt["cx"], top_corridor),
                   (tgt["cx"], tgt["top_y"])]

        # Pin Yes/No label at the actual exit vertex of the arrow.
        # The first waypoint (wps[0]) tells us which vertex the flow leaves from:
        #   right vertex  → (right_x, cy)   — forward flows
        #   bottom vertex → (cx, bot_y)     — skip-forward flows
        #   top vertex    → (cx, top_y)     — backward-loop flows
        # We match with a 2 px tolerance and place the label just outside that vertex.
        label_bounds = None
        if name and "Gateway" in src["type"] and src["type"] != "parallelGateway":
            ex, ey = wps[0]
            if abs(ex - src["right_x"]) <= 2:          # exits RIGHT → label above the line
                label_bounds = (ex + 4,  ey - 18, 30, 14)
            elif abs(ey - src["bot_y"]) <= 2:           # exits BOTTOM → label right of line
                label_bounds = (ex + 4,  ey + 4,  30, 14)
            elif abs(ey - src["top_y"]) <= 2:           # exits TOP (backward loop) → above
                label_bounds = (ex + 4,  ey - 18, 30, 14)
            else:                                        # fallback
                label_bounds = (ex + 4,  ey - 14, 30, 14)

        flows.append({"id": fid, "source": src_id, "target": tgt_id,
                      "name": name, "waypoints": wps, "label_bounds": label_bounds})

    for i in range(len(elements) - 1):
        src = elements[i]
        nxt = elements[i + 1]
        if "Gateway" in src["type"]:
            is_parallel = src["type"] == "parallelGateway"
            # Always use "Yes"/"No" — short labels that never overlap adjacent shapes
            add_flow(src["id"], nxt["id"], "" if is_parallel else "Yes")
            no_tgt = gateway_no.get(src["id"]) or "End_1"
            # Fallback 1: LLM gave a non-existent step ID — use End_1.
            if no_tgt not in el_by_id:
                no_tgt = "End_1"
            # Fallback 2: LLM set no_to == the Yes-path target (same element).
            # Both connectors would share an identical route and render as ONE
            # overlapping arrow with both labels on it.  Force No to End_1.
            if no_tgt == nxt["id"]:
                no_tgt = "End_1"
            # Only emit the No flow when it genuinely reaches a different element
            # (edge-case guard: if the gateway sits right before End_1, skip to
            # avoid two identical arrows — the diamond degrades to one exit).
            if no_tgt != nxt["id"]:
                add_flow(src["id"], no_tgt, "" if is_parallel else "No")
        else:
            add_flow(src["id"], nxt["id"])

    # ── Build XML ─────────────────────────────────────────────────────────────
    lines: list[str] = []
    lines.append('<?xml version="1.0" encoding="UTF-8"?>')
    lines.append('<definitions xmlns="http://www.omg.org/spec/BPMN/20100524/MODEL"')
    lines.append('             xmlns:bpmndi="http://www.omg.org/spec/BPMN/20100524/DI"')
    lines.append('             xmlns:dc="http://www.omg.org/spec/DD/20100524/DC"')
    lines.append('             xmlns:di="http://www.omg.org/spec/DD/20100524/DI"')
    lines.append('             id="Definitions_1" targetNamespace="http://bpmn.io/schema/bpmn">')

    # Collaboration wrapper (required for BPMN 2.0 pools with lanes)
    lines.append('  <collaboration id="Collab_1">')
    lines.append(f'    <participant id="Part_1" name="{_esc(process_name)}" processRef="Process_1"/>')
    lines.append('  </collaboration>')

    # Process
    lines.append(f'  <process id="Process_1" name="{_esc(process_name)}" isExecutable="false">')

    # Lane set — one lane per unique role
    lines.append('    <laneSet id="LaneSet_1">')
    for li, role in enumerate(role_order):
        lid = f"Lane_{li + 1}"
        lines.append(f'      <lane id="{lid}" name="{_esc(role)}">')
        for el in elements:
            if el["lane_idx"] == li:
                lines.append(f'        <flowNodeRef>{el["id"]}</flowNodeRef>')
        lines.append(f'      </lane>')
    lines.append('    </laneSet>')

    # Flow nodes — emit typed BPMN 2.0 elements
    for el in elements:
        eid, ename, etype = el["id"], _esc(el["name"]), el["type"]
        if etype == "startEvent":
            lines.append(f'    <startEvent id="{eid}" name="{ename}"/>')
        elif etype == "endEvent":
            lines.append(f'    <endEvent id="{eid}" name="{ename}"/>')
        elif "Gateway" in etype:
            lines.append(f'    <{etype} id="{eid}" name="{ename}"/>')
        elif etype in TASK_TYPES:
            lines.append(f'    <{etype} id="{eid}" name="{ename}"/>')
        else:
            lines.append(f'    <userTask id="{eid}" name="{ename}"/>')

    # Sequence flows
    for fl in flows:
        na = f' name="{_esc(fl["name"])}"' if fl["name"] else ""
        lines.append(f'    <sequenceFlow id="{fl["id"]}" sourceRef="{fl["source"]}" targetRef="{fl["target"]}"{na}/>')

    lines.append('  </process>')

    # ── Diagram section ───────────────────────────────────────────────────────
    lines.append('  <bpmndi:BPMNDiagram id="BPMNDiagram_1">')
    lines.append('    <bpmndi:BPMNPlane id="BPMNPlane_1" bpmnElement="Collab_1">')

    # Participant (pool) shape
    lines.append(f'      <bpmndi:BPMNShape id="Shape_Part_1" bpmnElement="Part_1" isHorizontal="true">')
    lines.append(f'        <dc:Bounds x="{P_X}" y="{P_Y}" width="{P_W}" height="{P_H}"/>')
    lines.append(f'      </bpmndi:BPMNShape>')

    # Lane shapes
    for li, role in enumerate(role_order):
        lid = f"Lane_{li + 1}"
        lines.append(f'      <bpmndi:BPMNShape id="Shape_{lid}" bpmnElement="{lid}" isHorizontal="true">')
        lines.append(f'        <dc:Bounds x="{P_X + POOL_HDR_W}" y="{P_Y + li * LANE_H}" width="{P_W - POOL_HDR_W}" height="{LANE_H}"/>')
        lines.append(f'      </bpmndi:BPMNShape>')

    # Element shapes
    _TYPED_TASKS = {"userTask", "serviceTask", "manualTask", "businessRuleTask",
                    "sendTask", "receiveTask", "scriptTask"}
    for el in elements:
        lines.append(f'      <bpmndi:BPMNShape id="Shape_{el["id"]}" bpmnElement="{el["id"]}">')
        lines.append(f'        <dc:Bounds x="{el["x"]}" y="{el["y"]}" width="{el["w"]}" height="{el["h"]}"/>')
        if "Gateway" in el["type"]:
            # Place the question label to the TOP-LEFT of the diamond so it
            # never overlaps the Yes label (right exit) or No label (bottom exit).
            # Right-align the label box at the diamond centre — text stays
            # entirely left of the Yes arrow which exits from the right vertex.
            # With H_STEP=175 and task_w=120 there is ~90 px of clear space
            # between the incoming task's right edge and the diamond centre,
            # so lx = cx - GW_LBL_W lands safely in that gap.
            GW_LBL_W = 90
            GW_LBL_H = 40
            GAP       = 6
            lx = el["cx"] - GW_LBL_W   # right edge at diamond centre
            ly = el["top_y"] - GW_LBL_H - GAP
            lines.append(f'        <bpmndi:BPMNLabel>')
            lines.append(f'          <dc:Bounds x="{lx}" y="{ly}" width="{GW_LBL_W}" height="{GW_LBL_H}"/>')
            lines.append(f'        </bpmndi:BPMNLabel>')
        elif el["type"] in _TYPED_TASKS:
            # bpmn-js renders the task-type marker icon at the top-left corner of
            # the shape (5 px inset, ~15 px tall → bottom edge ≈ shape.y + 20).
            # Set explicit label bounds that start BELOW the icon so the text
            # never overlaps it.  Height is reduced accordingly.
            ICON_H = 24
            lines.append(f'        <bpmndi:BPMNLabel>')
            lines.append(f'          <dc:Bounds x="{el["x"]}" y="{el["y"] + ICON_H}" width="{el["w"]}" height="{el["h"] - ICON_H}"/>')
            lines.append(f'        </bpmndi:BPMNLabel>')
        lines.append(f'      </bpmndi:BPMNShape>')

    # Edge shapes
    for fl in flows:
        if not fl["waypoints"]:
            continue
        lines.append(f'      <bpmndi:BPMNEdge id="Edge_{fl["id"]}" bpmnElement="{fl["id"]}">')
        for wx, wy in fl["waypoints"]:
            lines.append(f'        <di:waypoint x="{wx}" y="{wy}"/>')
        if fl.get("label_bounds"):
            lx, ly, lw, lh = fl["label_bounds"]
            lines.append(f'        <bpmndi:BPMNLabel>')
            lines.append(f'          <dc:Bounds x="{lx}" y="{ly}" width="{lw}" height="{lh}"/>')
            lines.append(f'        </bpmndi:BPMNLabel>')
        lines.append(f'      </bpmndi:BPMNEdge>')

    lines.append('    </bpmndi:BPMNPlane>')
    lines.append('  </bpmndi:BPMNDiagram>')
    lines.append('</definitions>')

    return "\n".join(lines)


def _esc(text: str) -> str:
    """Escape XML special characters in attribute values."""
    return (text
            .replace("&", "&amp;")
            .replace('"', "&quot;")
            .replace("<", "&lt;")
            .replace(">", "&gt;"))


def validate_xml(xml_str: str) -> tuple[bool, str]:
    try:
        ET.fromstring(xml_str.encode("utf-8"))
        return True, ""
    except ET.ParseError as e:
        return False, str(e)


# ── BPMN parser ───────────────────────────────────────────────────────────────

def _parse_bpmn(bpmn_xml: str):
    """Parse BPMN 2.0 XML from bpmn-js.

    Returns (proc_shapes, cont_shapes, connections) where:
      proc_shapes  = process elements (tasks/events/gateways):
                     [{'id','type','name','x','y','w','h'}, ...]  pixels
      cont_shapes  = container elements (participants/lanes):
                     [{'id','type','name','x','y','w','h'}, ...]  pixels
      connections  = sequence flows:
                     [{'id','source','target','wps','name'}, ...]

    Returns ([], [], []) on any error so callers can fall back to PNG.

    NOTE: cont_shapes are ordered participants-first so renderers can draw
    the pool border before the individual lane backgrounds.
    """
    if not bpmn_xml:
        return [], [], []
    # Types that are pool/lane containers — NOT rendered as process shapes
    CONTAINER_TYPES = {'participant', 'lane', 'collaboration', 'process', 'laneset'}
    try:
        root = ET.fromstring(bpmn_xml)
        BPMNDI = 'http://www.omg.org/spec/BPMN/20100524/DI'
        DC     = 'http://www.omg.org/spec/DD/20100524/DC'
        DI     = 'http://www.omg.org/spec/DD/20100524/DI'

        sem: dict = {}
        for el in root.iter():
            eid = el.get('id')
            if eid:
                tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
                sem[eid] = {
                    'type':   tag,
                    'name':   el.get('name', ''),
                    'source': el.get('sourceRef', ''),
                    'target': el.get('targetRef', ''),
                }

        proc_shapes: list = []
        cont_shapes: list = []
        for s in root.iter(f'{{{BPMNDI}}}BPMNShape'):
            eid    = s.get('bpmnElement', '')
            bounds = s.find(f'{{{DC}}}Bounds')
            if bounds is None:
                continue
            info = sem.get(eid, {'type': 'task', 'name': ''})
            entry = {
                'id':   eid,
                'type': info['type'],
                'name': info['name'],
                'x':    float(bounds.get('x', 0)),
                'y':    float(bounds.get('y', 0)),
                'w':    float(bounds.get('width',  100)),
                'h':    float(bounds.get('height',  80)),
            }
            if info['type'].lower() in CONTAINER_TYPES:
                cont_shapes.append(entry)
            else:
                proc_shapes.append(entry)

        # Participants first so renderers draw pool border before lane fills
        cont_shapes.sort(key=lambda s: (0 if s['type'].lower() == 'participant' else 1,
                                        s['y'], s['x']))

        connections: list = []
        for e in root.iter(f'{{{BPMNDI}}}BPMNEdge'):
            eid  = e.get('bpmnElement', '')
            info = sem.get(eid, {})
            wps  = [(float(wp.get('x', 0)), float(wp.get('y', 0)))
                    for wp in e.iter(f'{{{DI}}}waypoint')]
            connections.append({
                'id':     eid,
                'source': info.get('source', ''),
                'target': info.get('target', ''),
                'wps':    wps,
                'name':   info.get('name', ''),
            })
        return proc_shapes, cont_shapes, connections
    except Exception:
        return [], [], []


# ── Export helpers ────────────────────────────────────────────────────────────

def _make_shape_key_image() -> bytes:
    """Render a compact BPMN shape key as a PNG using PIL."""
    from PIL import Image, ImageDraw, ImageFont

    FONT_SIZE_TITLE = 13
    FONT_SIZE_LABEL = 11
    font_title = font_label = None
    for fpath in [
        "C:/Windows/Fonts/calibrib.ttf",
        "C:/Windows/Fonts/calibri.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ]:
        try:
            font_title = ImageFont.truetype(fpath, FONT_SIZE_TITLE)
            font_label = ImageFont.truetype(fpath, FONT_SIZE_LABEL)
            break
        except Exception:
            pass
    if font_title is None:
        font_title = ImageFont.load_default()
        font_label = font_title

    entries = [
        ("start",  "Start Event"),
        ("end",    "End Event"),
        ("xgw",    "Exclusive Gateway"),
        ("pgw",    "Parallel Gateway"),
        ("igw",    "Inclusive Gateway"),
        ("user",   "User Task"),
        ("svc",    "Service Task"),
        ("manual", "Manual Task"),
        ("br",     "Business Rule Task"),
        ("send",   "Send Task"),
        ("recv",   "Receive Task"),
    ]

    COLS     = 2
    COL_W    = 185
    ROW_H    = 24
    ICON_SZ  = 15
    ICON_PAD = 8
    TEXT_OFF = ICON_PAD + ICON_SZ + 6
    TITLE_H  = 22
    PAD_B    = 8

    rows = (len(entries) + COLS - 1) // COLS
    W = COLS * COL_W + 16
    H = TITLE_H + rows * ROW_H + PAD_B

    img = Image.new("RGBA", (W, H), (255, 255, 255, 255))
    d   = ImageDraw.Draw(img)

    # Border + title bar
    d.rectangle([0, 0, W - 1, H - 1], outline=(200, 200, 200), width=1)
    d.rectangle([0, 0, W - 1, TITLE_H - 1], fill=(248, 250, 252))
    d.line([(0, TITLE_H - 1), (W, TITLE_H - 1)], fill=(220, 220, 220))
    d.text((8, 4), "SHAPE KEY", fill=(100, 116, 139), font=font_title)

    # Column divider
    d.line([(COL_W + 8, TITLE_H), (COL_W + 8, H - 1)], fill=(230, 230, 230))

    for i, (key, label) in enumerate(entries):
        col = i % COLS
        row = i // COLS
        x0  = 8 + col * COL_W
        y0  = TITLE_H + row * ROW_H
        icx = x0 + ICON_PAD + ICON_SZ // 2
        icy = y0 + ROW_H // 2
        r   = ICON_SZ // 2 - 1

        if key == "start":
            d.ellipse([icx - r, icy - r, icx + r, icy + r],
                      outline=(34, 197, 94), fill=(255, 255, 255), width=2)
        elif key == "end":
            d.ellipse([icx - r, icy - r, icx + r, icy + r],
                      outline=(153, 27, 27), fill=(255, 255, 255), width=4)
        elif key in ("xgw", "pgw", "igw"):
            pts = [(icx, icy - r), (icx + r, icy), (icx, icy + r), (icx - r, icy)]
            d.polygon(pts, outline=(85, 85, 85), fill=(245, 203, 92))
            if key == "xgw":
                d.line([icx - 3, icy - 3, icx + 3, icy + 3], fill=(85, 85, 85), width=1)
                d.line([icx + 3, icy - 3, icx - 3, icy + 3], fill=(85, 85, 85), width=1)
            elif key == "pgw":
                d.line([icx, icy - 3, icx, icy + 3], fill=(85, 85, 85), width=1)
                d.line([icx - 3, icy, icx + 3, icy], fill=(85, 85, 85), width=1)
            else:
                d.ellipse([icx - 2, icy - 2, icx + 2, icy + 2],
                          outline=(85, 85, 85), fill=(245, 203, 92), width=1)
        else:
            rw, rh = ICON_SZ, ICON_SZ - 2
            d.rectangle([icx - rw // 2, icy - rh // 2, icx + rw // 2, icy + rh // 2],
                        outline=(100, 130, 160), fill=(175, 193, 214))
            if key == "user":
                d.ellipse([icx - 2, icy - rh // 2 + 1, icx + 2, icy - rh // 2 + 5],
                          fill=(60, 80, 100))
                d.arc([icx - 4, icy - 2, icx + 4, icy + rh // 2 - 1],
                      start=0, end=180, fill=(60, 80, 100), width=1)
            elif key == "svc":
                d.ellipse([icx - 2, icy - 2, icx + 2, icy + 2],
                          outline=(60, 80, 100), fill=(175, 193, 214), width=1)
                for ang_d in [0, 90, 45, 135]:
                    import math
                    rad = math.radians(ang_d)
                    d.line([icx + int(4 * math.cos(rad)), icy + int(4 * math.sin(rad)),
                            icx - int(4 * math.cos(rad)), icy - int(4 * math.sin(rad))],
                           fill=(60, 80, 100), width=1)
            elif key == "br":
                d.rectangle([icx - 4, icy - 3, icx + 4, icy - 1], fill=(60, 80, 100))
                d.line([icx - 4, icy + 1, icx + 4, icy + 1], fill=(120, 140, 160), width=1)
                d.line([icx - 4, icy + 3, icx + 4, icy + 3], fill=(160, 175, 190), width=1)
            elif key == "send":
                d.rectangle([icx - 4, icy - 2, icx + 4, icy + 3], fill=(60, 80, 100))
                d.line([icx - 4, icy - 2, icx, icy + 1], fill=(255, 255, 255), width=1)
                d.line([icx + 4, icy - 2, icx, icy + 1], fill=(255, 255, 255), width=1)
            elif key == "recv":
                d.rectangle([icx - 4, icy - 2, icx + 4, icy + 3],
                            outline=(60, 80, 100), fill=(175, 193, 214), width=1)
                d.line([icx - 4, icy - 2, icx, icy + 1], fill=(60, 80, 100), width=1)
                d.line([icx + 4, icy - 2, icx, icy + 1], fill=(60, 80, 100), width=1)

        d.text((x0 + TEXT_OFF, icy - 6), label, fill=(60, 80, 100), font=font_label)

    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


def _make_pptx(png_bytes: bytes, process_name: str, bpmn_xml: str = '') -> bytes:
    from pptx import Presentation as PRS
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

    prs = PRS()
    proc_shapes, cont_shapes, connections = _parse_bpmn(bpmn_xml)

    if proc_shapes:
        # ── coordinate mapping ───────────────────────────────────────────────
        # Use the full extent (including lane/pool containers) so shape positions
        # are consistent with the container backgrounds.
        all_shp = proc_shapes + cont_shapes
        min_x = min(s['x']           for s in all_shp)
        min_y = min(s['y']           for s in all_shp)
        max_x = max(s['x'] + s['w'] for s in all_shp)
        max_y = max(s['y'] + s['h'] for s in all_shp)

        MARGIN  = 304800          # 1/3" in EMU
        TITLE_H = 457200          # 0.5" in EMU
        src_w   = max(max_x - min_x, 1)
        src_h   = max(max_y - min_y, 1)

        # Natural scale: 1 BPMN pixel = 1/96 inch = 9525 EMU.
        # This keeps shapes at their natural readable size (tasks ~1.25" wide).
        # The slide expands to fit the diagram rather than squishing it.
        SCALE = 914400 // 96   # 9525 EMU/px

        # Set slide dimensions to fit the full diagram at natural scale
        prs.slide_width  = int(src_w * SCALE + 2 * MARGIN)
        prs.slide_height = int(src_h * SCALE + TITLE_H + 2 * MARGIN)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        def bx(px): return MARGIN + int((px - min_x) * SCALE)
        def by(py): return MARGIN + TITLE_H + int((py - min_y) * SCALE)
        def bw(pw): return max(int(pw * SCALE), 1)

        # ── title ─────────────────────────────────────────────────────────────
        title_w = int(src_w * SCALE)  # span full diagram width
        tb = slide.shapes.add_textbox(MARGIN, int(MARGIN * 0.4), title_w, TITLE_H)
        tf = tb.text_frame
        tf.text = process_name
        run = tf.paragraphs[0].runs[0]
        run.font.size  = Pt(16)
        run.font.bold  = True
        run.font.color.rgb = RGBColor(0x1F, 0x2D, 0x3D)

        # ── 1. Lane / pool backgrounds (drawn first → bottom of z-stack) ─────
        LANE_FILLS   = [RGBColor(0xDA, 0xE8, 0xFC), RGBColor(0xEB, 0xF3, 0xFF)]
        LANE_BORDER  = RGBColor(0x6C, 0x8E, 0xBF)
        PART_BORDER  = RGBColor(0x4A, 0x6E, 0x9A)
        lanes = [s for s in cont_shapes if s['type'].lower() == 'lane']
        parts = [s for s in cont_shapes if s['type'].lower() == 'participant']

        for s in parts:
            sx, sy, sw, sh = bx(s['x']), by(s['y']), bw(s['w']), bw(s['h'])
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, sx, sy, sw, sh)
            rect.fill.background()              # transparent interior
            rect.line.color.rgb = PART_BORDER
            rect.line.width     = Emu(19050)    # 1.5 pt

        for li, s in enumerate(lanes):
            sx, sy, sw, sh = bx(s['x']), by(s['y']), bw(s['w']), bw(s['h'])
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, sx, sy, sw, sh)
            rect.fill.solid()
            rect.fill.fore_color.rgb = LANE_FILLS[li % 2]
            rect.line.color.rgb      = LANE_BORDER
            rect.line.width          = Emu(9525)
            # No text in the lane rectangle itself — use a separate label text box
            # so the lane name stays in a controlled area at the left of the lane.
            if s['name']:
                # Narrow textbox at the left of the lane — width matches
                # LANE_LBL_W (120 px) so text is always clear of process shapes.
                lbl_w = bw(120)
                lbl   = slide.shapes.add_textbox(sx, sy, lbl_w, sh)
                lf    = lbl.text_frame
                lf.word_wrap = True
                lf.auto_size = MSO_AUTO_SIZE.NONE
                lf.text = s['name']
                for para in lf.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
                    for r2 in para.runs:
                        r2.font.size      = Pt(8)
                        r2.font.bold      = True
                        r2.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

        # ── 2. Connectors (behind process shapes) ─────────────────────────────
        from pptx.oxml import parse_xml as _pptx_parse_xml
        _TAIL_ARROW = (
            '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            ' type="arrow" w="med" len="med"/>'
        )
        for conn in connections:
            wps        = conn['wps']
            flow_label = conn.get('name', '')
            if len(wps) < 2:
                continue
            # Draw each waypoint segment so bent/multi-hop paths render correctly.
            # Only the LAST segment gets the arrowhead (pointing at the target shape).
            for seg_i in range(len(wps) - 1):
                x0, y0 = bx(wps[seg_i][0]),     by(wps[seg_i][1])
                x1, y1 = bx(wps[seg_i + 1][0]), by(wps[seg_i + 1][1])
                cxn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x0, y0, x1, y1)
                cxn.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
                cxn.line.width     = Emu(12700)    # 1 pt
                if seg_i == len(wps) - 2:          # last segment → add arrowhead
                    cxn.line._ln.append(_pptx_parse_xml(_TAIL_ARROW))

            # ── Yes / No (or other flow name) label near source exit point ──────
            # Position depends on whether the flow exits horizontally (→ label
            # above the line) or vertically (→ label to the right of the line).
            if flow_label:
                ex_px, ey_px = wps[0]
                d_x = wps[1][0] - wps[0][0]
                d_y = wps[1][1] - wps[0][1]
                if abs(d_x) >= abs(d_y):    # horizontal exit: label above
                    lx_px, ly_px = ex_px + 4, ey_px - 16
                else:                        # vertical exit: label right of line
                    lx_px, ly_px = ex_px + 4, ey_px + 4
                flbl = slide.shapes.add_textbox(
                    bx(lx_px), by(ly_px), bw(35), bw(14))
                flf  = flbl.text_frame
                flf.word_wrap = False
                flf.text      = flow_label
                for para in flf.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
                    for run in para.runs:
                        run.font.size      = Pt(8)
                        run.font.italic    = True
                        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

        # ── 3. Process shapes (top of z-stack) ────────────────────────────────
        # Matching bpmn-js rendering exactly:
        #   • Tasks/subprocesses → text INSIDE the shape (auto-shrink to fit)
        #   • Events             → oval with no text, label text box BELOW
        #   • Gateways           → diamond with no text, label text box ABOVE
        SHAPE_CFG = {
            'startevent':        (MSO_AUTO_SHAPE_TYPE.OVAL,            RGBColor(0x67, 0xAB, 0x9F)),
            'endevent':          (MSO_AUTO_SHAPE_TYPE.OVAL,            RGBColor(0xC0, 0x39, 0x2B)),
            'intermediateevent': (MSO_AUTO_SHAPE_TYPE.OVAL,            RGBColor(0xF3, 0x9C, 0x12)),
            'exclusivegateway':  (MSO_AUTO_SHAPE_TYPE.DIAMOND,         RGBColor(0xF5, 0xCB, 0x5C)),
            'inclusivegateway':  (MSO_AUTO_SHAPE_TYPE.DIAMOND,         RGBColor(0xF5, 0xCB, 0x5C)),
            'parallelgateway':   (MSO_AUTO_SHAPE_TYPE.DIAMOND,         RGBColor(0xF5, 0xCB, 0x5C)),
            'subprocess':        (MSO_AUTO_SHAPE_TYPE.RECTANGLE,       RGBColor(0xAF, 0xC1, 0xD6)),
        }
        DEFAULT_CFG = (MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, RGBColor(0xDA, 0xE8, 0xFC))

        # Label box dimensions matching build_bpmn_xml gateway label bounds.
        # Width=90 px, Height=50 px (fits 3 lines). Anchored so the box BOTTOM
        # is 4 px above the diamond top — same logic as the BPMN XML label.
        LBL_W_PX = 90
        LBL_H_PX = 50   # increased from 28 so multi-line labels don't overflow
        GAP_PX   = 4

        for s in proc_shapes:
            typ = s['type'].lower()
            key = next((k for k in SHAPE_CFG if k in typ), None)
            stype, colour = SHAPE_CFG[key] if key else DEFAULT_CFG

            sx, sy, sw, sh = bx(s['x']), by(s['y']), bw(s['w']), bw(s['h'])
            pshape = slide.shapes.add_shape(stype, sx, sy, sw, sh)
            pshape.fill.solid()
            pshape.fill.fore_color.rgb = colour
            pshape.line.color.rgb      = RGBColor(0x55, 0x55, 0x55)
            pshape.line.width          = Emu(9525)

            is_event   = 'event'   in typ
            is_gateway = 'gateway' in typ

            if is_event or is_gateway:
                # Shape has no text — label goes in a separate text box
                pshape.text_frame.text = ''
                if s['name']:
                    # Events: label below the shape;  Gateways: label above.
                    # Gateway bottom-anchored: box bottom = shape top - GAP_PX
                    lbl_cx = s['x'] + s['w'] / 2 - LBL_W_PX / 2
                    if is_gateway:
                        lbl_y = s['y'] - LBL_H_PX - GAP_PX  # bottom edge = shape top - gap
                    else:
                        lbl_y = s['y'] + s['h'] + 2           # below event
                    lbl_sx = bx(lbl_cx)
                    lbl_sy = max(by(lbl_y), MARGIN + TITLE_H)
                    lbl_sw = bw(LBL_W_PX)
                    lbl_sh = bw(LBL_H_PX)
                    lbl = slide.shapes.add_textbox(lbl_sx, lbl_sy, lbl_sw, lbl_sh)
                    lf  = lbl.text_frame
                    lf.word_wrap = True
                    # Let the text box grow to fit the label text
                    lf.text = s['name']
                    for para in lf.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for r2 in para.runs:
                            r2.font.size      = Pt(8)
                            r2.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
            else:
                # Tasks / subprocesses: text inside, auto-shrink to fit shape
                tf = pshape.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                if s['name']:
                    tf.text = s['name']
                    for para in tf.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for r2 in para.runs:
                            r2.font.size      = Pt(8)
                            r2.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    else:
        # ── fallback: PNG image ────────────────────────────────────────────────
        prs.slide_width  = Emu(12192000)   # 13.33"
        prs.slide_height = Emu(6858000)    # 7.5"
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.6))
        tf = tb.text_frame
        tf.text = process_name
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(20); run.font.bold = True
        slide.shapes.add_picture(io.BytesIO(png_bytes),
                                 Inches(0.4), Inches(0.9), Inches(12.5), Inches(6.3))

    # ── Shape key: bottom-right corner ────────────────────────────────────────
    try:
        key_png    = _make_shape_key_image()
        KEY_W_E    = Emu(int(2.6 * 914400))
        KEY_H_E    = Emu(int(1.75 * 914400))
        _key_mar   = Emu(304800)   # 1/3" margin
        key_left   = prs.slide_width  - KEY_W_E - _key_mar
        key_top    = prs.slide_height - KEY_H_E - _key_mar
        slide.shapes.add_picture(io.BytesIO(key_png), key_left, key_top, KEY_W_E, KEY_H_E)
    except Exception:
        pass

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _make_docx(png_bytes: bytes, process_name: str, bpmn_xml: str = '') -> bytes:
    from docx import Document
    from docx.shared import Inches, Pt
    from lxml import etree

    doc = Document()
    sec = doc.sections[0]

    proc_shapes, cont_shapes, connections = _parse_bpmn(bpmn_xml)

    if proc_shapes:
        DPI       = 96.0
        EMU_IN    = 914400        # EMU per inch
        MARGIN_IN = 0.5
        TITLE_IN  = 0.45

        all_shp = proc_shapes + cont_shapes
        min_x = min(s['x']           for s in all_shp)
        min_y = min(s['y']           for s in all_shp)
        max_x = max(s['x'] + s['w'] for s in all_shp)
        max_y = max(s['y'] + s['h'] for s in all_shp)

        # OOXML page size limit: 31 680 twips = 22 inches per dimension.
        # If the diagram at natural 1px=1/96" scale exceeds 22" wide, scale
        # it down proportionally so everything fits horizontally.
        # Page height is fixed to landscape-letter (11"); content taller than
        # that spills naturally to subsequent pages via the floating-shape y-offset.
        OOXML_MAX_IN  = 22.0    # Word/OOXML hard cap per axis
        PAGE_H_IN     = 11.0    # landscape-letter height (standard print size)

        src_w_in = max((max_x - min_x) / DPI, 0.1)
        src_h_in = max((max_y - min_y) / DPI, 0.1)

        usable_w_in = OOXML_MAX_IN - 2 * MARGIN_IN
        # Scale = 1.0 (natural) unless diagram is wider than the max usable width
        scale = min(usable_w_in / src_w_in, 1.0)

        page_w_in = src_w_in * scale + 2 * MARGIN_IN   # always ≤ OOXML_MAX_IN
        # Height: use the standard page height; tall diagrams overflow to page 2+
        page_h_in = PAGE_H_IN

        sec.page_width   = int(page_w_in * EMU_IN)
        sec.page_height  = int(page_h_in * EMU_IN)
        sec.left_margin  = sec.right_margin  = int(MARGIN_IN * EMU_IN)
        sec.top_margin   = sec.bottom_margin = int(MARGIN_IN * EMU_IN)

        def ex(px): return int((MARGIN_IN + (px - min_x) / DPI * scale) * EMU_IN)
        def ey(py): return int((MARGIN_IN + TITLE_IN + (py - min_y) / DPI * scale) * EMU_IN)
        def ew(pw): return max(int(pw / DPI * scale * EMU_IN), 91440)

        FILL = {
            'startevent': '67AB9F', 'endevent': 'C0392B',
            'intermediateevent': 'F39C12',
            'exclusivegateway': 'F5CB5C', 'inclusivegateway': 'F5CB5C',
            'parallelgateway': 'F5CB5C', 'subprocess': 'AFC1D6',
        }
        DEFAULT_FILL = 'DAE8FC'
        LANE_FILLS   = ['DAE8FC', 'EBF3FF']
        GEOM = {'event': 'ellipse', 'gateway': 'diamond'}

        W   = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
        WP  = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
        A   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        WPS = 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape'

        # Use a single empty paragraph as the anchor for all floating shapes.
        # The title is rendered as a floating text box so the custom page
        # dimensions fully control layout without the heading paragraph
        # causing Word to extend the visible area.
        float_para = doc.add_paragraph()
        run = float_para.add_run()
        sid = 1

        # ── 0. Title floating text box ─────────────────────────────────────
        title_x = int(MARGIN_IN * EMU_IN)
        title_y = int(MARGIN_IN * 0.3 * EMU_IN)
        title_w = int(src_w_in * scale * EMU_IN)
        title_h = int(TITLE_IN * EMU_IN)
        title_xml = (
            f'<w:drawing xmlns:w="{W}">'
            f'<wp:anchor distT="0" distB="0" distL="0" distR="0" simplePos="0"'
            f' relativeHeight="500" behindDoc="0" locked="0"'
            f' layoutInCell="1" allowOverlap="1" xmlns:wp="{WP}">'
            f'<wp:simplePos x="0" y="0"/>'
            f'<wp:positionH relativeFrom="page"><wp:posOffset>{title_x}</wp:posOffset></wp:positionH>'
            f'<wp:positionV relativeFrom="page"><wp:posOffset>{title_y}</wp:posOffset></wp:positionV>'
            f'<wp:extent cx="{title_w}" cy="{title_h}"/>'
            f'<wp:effectExtent l="0" t="0" r="0" b="0"/>'
            f'<wp:wrapNone/>'
            f'<wp:docPr id="0" name="Title"/>'
            f'<wp:cNvGraphicFramePr/>'
            f'<a:graphic xmlns:a="{A}">'
            f'<a:graphicData uri="{WPS}">'
            f'<wps:wsp xmlns:wps="{WPS}">'
            f'<wps:cNvSpPr><a:spLocks noChangeArrowheads="1"/></wps:cNvSpPr>'
            f'<wps:spPr>'
            f'<a:xfrm><a:off x="0" y="0"/><a:ext cx="{title_w}" cy="{title_h}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:noFill/><a:ln><a:noFill/></a:ln>'
            f'</wps:spPr>'
            f'<wps:txbx>'
            f'<w:txbxContent>'
            f'<w:p><w:pPr><w:jc w:val="left"/></w:pPr>'
            f'<w:r><w:rPr><w:b/><w:sz w:val="28"/><w:szCs w:val="28"/>'
            f'<w:color w:val="1F2D3D"/></w:rPr>'
            f'<w:t xml:space="preserve">{_esc(process_name)}</w:t></w:r></w:p>'
            f'</w:txbxContent>'
            f'</wps:txbx>'
            f'<wps:bodyPr anchor="ctr"/>'
            f'</wps:wsp>'
            f'</a:graphicData>'
            f'</a:graphic>'
            f'</wp:anchor>'
            f'</w:drawing>'
        )
        run._r.append(etree.fromstring(title_xml.encode('utf-8')))

        def _anchor(sid, x_e, y_e, w_e, h_e, prst, fill_hex, label,
                    behind=False, rh=20000, text_align='center',
                    is_line=False, fh='0', fv='0', arrow=False,
                    no_fill=False):
            """Build a <w:drawing><wp:anchor> XML string for a floating shape.

            no_fill=True → transparent background and no border (used for
            external label text boxes placed above/below events and gateways).
            """
            bd  = '1' if behind else '0'
            if is_line:
                tail_xml = ('<a:tailEnd type="arrow" w="sm" len="sm"/>'
                            if arrow else '')
                spPr = (
                    f'<a:xfrm flipH="{fh}" flipV="{fv}">'
                    f'<a:off x="0" y="0"/><a:ext cx="{w_e}" cy="{h_e}"/></a:xfrm>'
                    f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
                    f'<a:noFill/>'
                    f'<a:ln w="9525"><a:solidFill><a:srgbClr val="444444"/></a:solidFill>'
                    f'{tail_xml}</a:ln>'
                )
                body   = '<wps:bodyPr/>'
                txbx   = ''
                cNvSp  = f'<wps:cNvSpPr isConnector="1"><a:spLocks noChangeArrowheads="1"/></wps:cNvSpPr>'
            else:
                if no_fill:
                    fill_ln = '<a:noFill/><a:ln><a:noFill/></a:ln>'
                else:
                    lnFill   = '6C8EBF' if fill_hex in (LANE_FILLS + ['FFFFFF']) else '555555'
                    fill_ln  = (f'<a:solidFill><a:srgbClr val="{fill_hex}"/></a:solidFill>'
                                f'<a:ln w="9525"><a:solidFill><a:srgbClr val="{lnFill}"/></a:solidFill></a:ln>')
                spPr = (
                    f'<a:xfrm><a:off x="0" y="0"/><a:ext cx="{w_e}" cy="{h_e}"/></a:xfrm>'
                    f'<a:prstGeom prst="{prst}"><a:avLst/></a:prstGeom>'
                    f'{fill_ln}'
                )
                esc_label = _esc(label) if label else ''
                txbx = (
                    f'<wps:txbx>'
                    f'<w:txbxContent>'
                    f'<w:p><w:pPr><w:jc w:val="{text_align}"/></w:pPr>'
                    f'<w:r><w:rPr><w:sz w:val="14"/><w:szCs w:val="14"/></w:rPr>'
                    f'<w:t xml:space="preserve">{esc_label}</w:t></w:r></w:p>'
                    f'</w:txbxContent>'
                    f'</wps:txbx>'
                )
                # normAutofit = "shrink text on overflow" (prevents text escaping shape)
                body  = '<wps:bodyPr anchor="ctr" lIns="45720" rIns="45720" tIns="36000" bIns="36000"><a:normAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/></wps:bodyPr>'
                cNvSp = f'<wps:cNvSpPr><a:spLocks noChangeArrowheads="1"/></wps:cNvSpPr>'
            return (
                f'<w:drawing xmlns:w="{W}">'
                f'<wp:anchor distT="0" distB="0" distL="0" distR="0" simplePos="0"'
                f' relativeHeight="{rh + sid}" behindDoc="{bd}" locked="0"'
                f' layoutInCell="1" allowOverlap="1" xmlns:wp="{WP}">'
                f'<wp:simplePos x="0" y="0"/>'
                f'<wp:positionH relativeFrom="page"><wp:posOffset>{x_e}</wp:posOffset></wp:positionH>'
                f'<wp:positionV relativeFrom="page"><wp:posOffset>{y_e}</wp:posOffset></wp:positionV>'
                f'<wp:extent cx="{w_e}" cy="{h_e}"/>'
                f'<wp:effectExtent l="0" t="0" r="0" b="0"/>'
                f'<wp:wrapNone/>'
                f'<wp:docPr id="{sid}" name="S{sid}"/>'
                f'<wp:cNvGraphicFramePr/>'
                f'<a:graphic xmlns:a="{A}">'
                f'<a:graphicData uri="{WPS}">'
                f'<wps:wsp xmlns:wps="{WPS}">'
                f'{cNvSp}'
                f'<wps:spPr>{spPr}</wps:spPr>'
                f'{txbx}{body}'
                f'</wps:wsp>'
                f'</a:graphicData>'
                f'</a:graphic>'
                f'</wp:anchor>'
                f'</w:drawing>'
            )

        # ── 1. Lane / pool backgrounds (behindDoc=1, low relativeHeight) ──────
        # Lane actor names are rendered in a SEPARATE narrow label box (120 px
        # wide) rather than inside the full-width lane rectangle. This prevents
        # the actor name text from extending into the start-event circle area.
        LANE_LBL_W_PX = 120   # must match LANE_LBL_W in build_bpmn_xml
        lanes_ordered = [s for s in cont_shapes if s['type'].lower() == 'lane']
        for s in cont_shapes:
            typ = s['type'].lower()
            x_e = ex(s['x']); y_e = ey(s['y'])
            w_e = ew(s['w']); h_e = ew(s['h'])
            if typ == 'participant':
                fill_hex = 'FFFFFF'
                rh       = 1000
                # Participant: render with pool name in full rect (no overlap risk)
                xml = _anchor(sid, x_e, y_e, w_e, h_e, 'rect', fill_hex,
                              s['name'], behind=True, rh=rh, text_align='left')
                run._r.append(etree.fromstring(xml.encode('utf-8')))
                sid += 1
            else:
                li       = lanes_ordered.index(s) if s in lanes_ordered else 0
                fill_hex = LANE_FILLS[li % 2]
                rh       = 2000 + li
                # Lane: render background rect with NO text to avoid overlap
                xml = _anchor(sid, x_e, y_e, w_e, h_e, 'rect', fill_hex,
                              '', behind=True, rh=rh, text_align='left')
                run._r.append(etree.fromstring(xml.encode('utf-8')))
                sid += 1
                # Separate narrow label box — actor name stays within 120 px,
                # always clear of the first process shape (start event).
                if s.get('name'):
                    lbl_w_e = ew(LANE_LBL_W_PX)
                    lbl_xml = _anchor(sid, x_e, y_e, lbl_w_e, h_e, 'rect',
                                      fill_hex, s['name'], behind=True,
                                      rh=rh + 1, text_align='left')
                    run._r.append(etree.fromstring(lbl_xml.encode('utf-8')))
                    sid += 1

        # ── 2. Connectors — one shape per waypoint segment ────────────────────
        # Drawing each segment individually ensures bent/routed BPMN paths
        # (L-shapes, corridor routes) render exactly as on the webpage.
        # Only the final segment carries the arrowhead.
        HALF_LW = 4763   # half of 9525 EMU (1 pt line) — used to centre
                         # near-zero-height/width shapes on the line axis
        for conn in connections:
            wps_pts = conn['wps']
            if len(wps_pts) < 2:
                continue
            n_segs = len(wps_pts) - 1
            for seg_i in range(n_segs):
                px1, py1 = wps_pts[seg_i]
                px2, py2 = wps_pts[seg_i + 1]
                x1, y1 = ex(px1), ey(py1)
                x2, y2 = ex(px2), ey(py2)
                dx, dy = x2 - x1, y2 - y1

                fh = '1' if dx < 0 else '0'
                fv = '1' if dy < 0 else '0'

                if abs(dx) < 100:          # essentially vertical segment
                    left_e = x1 - HALF_LW
                    top_e  = min(y1, y2)
                    w_seg  = 9525
                    h_seg  = max(abs(dy), 9525)
                elif abs(dy) < 100:        # essentially horizontal segment
                    left_e = min(x1, x2)
                    top_e  = y1 - HALF_LW
                    w_seg  = max(abs(dx), 9525)
                    h_seg  = 9525
                else:                      # diagonal segment
                    left_e = min(x1, x2)
                    top_e  = min(y1, y2)
                    w_seg  = abs(dx)
                    h_seg  = abs(dy)

                is_last = (seg_i == n_segs - 1)
                xml = _anchor(sid, left_e, top_e, w_seg, h_seg,
                              'line', '444444', '',
                              behind=False, rh=10000,
                              is_line=True, fh=fh, fv=fv, arrow=is_last)
                run._r.append(etree.fromstring(xml.encode('utf-8')))
                sid += 1

            # ── Yes / No (or other flow name) label near source exit ──────────
            flow_label = conn.get('name', '')
            if flow_label and len(wps_pts) >= 2:
                ex_px, ey_px = wps_pts[0]
                d_x = wps_pts[1][0] - wps_pts[0][0]
                d_y = wps_pts[1][1] - wps_pts[0][1]
                LBL_FW_PX = 30   # flow label box width (px)
                LBL_FH_PX = 14   # flow label box height (px)
                if abs(d_x) >= abs(d_y):   # horizontal exit → label above line
                    lx_px = ex_px + 4
                    ly_px = ey_px - 16
                else:                       # vertical exit → label right of line
                    lx_px = ex_px + 4
                    ly_px = ey_px + 4
                fl_x_e = ex(lx_px)
                fl_y_e = ey(ly_px)
                fl_w_e = ew(LBL_FW_PX)
                fl_h_e = ew(LBL_FH_PX)
                fl_xml = (
                    f'<w:drawing xmlns:w="{W}">'
                    f'<wp:anchor distT="0" distB="0" distL="0" distR="0" simplePos="0"'
                    f' relativeHeight="{10500 + sid}" behindDoc="0" locked="0"'
                    f' layoutInCell="1" allowOverlap="1" xmlns:wp="{WP}">'
                    f'<wp:simplePos x="0" y="0"/>'
                    f'<wp:positionH relativeFrom="page"><wp:posOffset>{fl_x_e}</wp:posOffset></wp:positionH>'
                    f'<wp:positionV relativeFrom="page"><wp:posOffset>{fl_y_e}</wp:posOffset></wp:positionV>'
                    f'<wp:extent cx="{fl_w_e}" cy="{fl_h_e}"/>'
                    f'<wp:effectExtent l="0" t="0" r="0" b="0"/>'
                    f'<wp:wrapNone/>'
                    f'<wp:docPr id="{sid}" name="FL{sid}"/>'
                    f'<wp:cNvGraphicFramePr/>'
                    f'<a:graphic xmlns:a="{A}">'
                    f'<a:graphicData uri="{WPS}">'
                    f'<wps:wsp xmlns:wps="{WPS}">'
                    f'<wps:cNvSpPr><a:spLocks noChangeArrowheads="1"/></wps:cNvSpPr>'
                    f'<wps:spPr>'
                    f'<a:xfrm><a:off x="0" y="0"/><a:ext cx="{fl_w_e}" cy="{fl_h_e}"/></a:xfrm>'
                    f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                    f'<a:noFill/><a:ln><a:noFill/></a:ln>'
                    f'</wps:spPr>'
                    f'<wps:txbx>'
                    f'<w:txbxContent>'
                    f'<w:p><w:pPr><w:jc w:val="left"/></w:pPr>'
                    f'<w:r><w:rPr><w:i/><w:sz w:val="12"/><w:szCs w:val="12"/>'
                    f'<w:color w:val="444444"/></w:rPr>'
                    f'<w:t xml:space="preserve">{_esc(flow_label)}</w:t></w:r></w:p>'
                    f'</w:txbxContent>'
                    f'</wps:txbx>'
                    f'<wps:bodyPr anchor="ctr"/>'
                    f'</wps:wsp>'
                    f'</a:graphicData>'
                    f'</a:graphic>'
                    f'</wp:anchor>'
                    f'</w:drawing>'
                )
                run._r.append(etree.fromstring(fl_xml.encode('utf-8')))
                sid += 1

        # ── 3. Process shapes (foreground) ────────────────────────────────────
        # Gateways and events follow the BPMN convention: the shape carries
        # NO text; instead a separate transparent text box is placed outside
        # (above gateway, below event) — matching the bpmn-js webpage rendering.
        # Label box: 90 px wide, 50 px tall (3-line capacity).
        # Gateway: bottom-anchored so box bottom = shape top - GAP_PX.
        # Event:   top-anchored just below the event circle.
        LBL_W_PX = 90
        LBL_H_PX = 50
        GAP_PX   = 4

        for s in proc_shapes:
            typ  = s['type'].lower()
            fkey = next((k for k in FILL if k in typ), None)
            fill = FILL[fkey] if fkey else DEFAULT_FILL
            gkey = next((k for k in GEOM if k in typ), None)
            prst = GEOM[gkey] if gkey else 'roundRect'
            x_e = ex(s['x']); y_e = ey(s['y'])
            w_e = ew(s['w']); h_e = ew(s['h'])

            is_event   = 'event'   in typ
            is_gateway = 'gateway' in typ

            # Shape itself — no text for events/gateways
            shape_label = '' if (is_event or is_gateway) else s['name']
            xml = _anchor(sid, x_e, y_e, w_e, h_e, prst, fill, shape_label,
                          behind=False, rh=20000)
            run._r.append(etree.fromstring(xml.encode('utf-8')))
            sid += 1

            # External label text box for events and gateways
            if (is_event or is_gateway) and s['name']:
                lbl_cx_px = s['x'] + s['w'] / 2 - LBL_W_PX / 2
                if is_gateway:
                    lbl_y_px = s['y'] - LBL_H_PX - GAP_PX  # bottom = shape top - gap
                else:
                    lbl_y_px = s['y'] + s['h'] + 2           # below event circle
                lbl_x_e = ex(lbl_cx_px)
                lbl_y_e = ey(lbl_y_px)
                lbl_w_e = ew(LBL_W_PX)
                lbl_h_e = ew(LBL_H_PX)
                xml = _anchor(sid, lbl_x_e, lbl_y_e, lbl_w_e, lbl_h_e,
                              'rect', None, s['name'],
                              behind=False, rh=20001, no_fill=True)
                run._r.append(etree.fromstring(xml.encode('utf-8')))
                sid += 1

    else:
        # ── fallback: embedded PNG ─────────────────────────────────────────
        sec.page_width   = Inches(11)
        sec.page_height  = Inches(8.5)
        sec.left_margin  = sec.right_margin  = Inches(0.5)
        sec.top_margin   = sec.bottom_margin = Inches(0.5)
        doc.add_heading(process_name, 0)
        doc.add_picture(io.BytesIO(png_bytes), width=Inches(9.5))

    # ── Shape key: floating picture at bottom-right ───────────────────────────
    try:
        key_png  = _make_shape_key_image()
        KEY_W_IN = 2.6
        KEY_H_IN = 1.75
        KEY_W_E  = int(KEY_W_IN * EMU_IN)
        KEY_H_E  = int(KEY_H_IN * EMU_IN)

        # Add image as a document relationship via a temporary inline picture
        tmp_p  = doc.add_paragraph()
        tmp_r  = tmp_p.add_run()
        tmp_il = tmp_r.add_picture(io.BytesIO(key_png), width=Inches(KEY_W_IN))
        blip_el = tmp_il._inline.find(
            './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
        r_id = blip_el.get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
        tmp_p._p.getparent().remove(tmp_p._p)

        # Position: bottom-right of the content area
        key_x_e = int((page_w_in - MARGIN_IN - KEY_W_IN) * EMU_IN)
        key_y_e = int((MARGIN_IN + TITLE_IN + src_h_in * scale - KEY_H_IN) * EMU_IN)

        PIC_NS = 'http://schemas.openxmlformats.org/drawingml/2006/picture'
        R_NS   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        pic_xml = (
            f'<w:drawing xmlns:w="{W}">'
            f'<wp:anchor distT="0" distB="0" distL="0" distR="0" simplePos="0"'
            f' relativeHeight="30000" behindDoc="0" locked="0"'
            f' layoutInCell="1" allowOverlap="1" xmlns:wp="{WP}">'
            f'<wp:simplePos x="0" y="0"/>'
            f'<wp:positionH relativeFrom="page"><wp:posOffset>{key_x_e}</wp:posOffset></wp:positionH>'
            f'<wp:positionV relativeFrom="page"><wp:posOffset>{key_y_e}</wp:posOffset></wp:positionV>'
            f'<wp:extent cx="{KEY_W_E}" cy="{KEY_H_E}"/>'
            f'<wp:effectExtent l="0" t="0" r="0" b="0"/>'
            f'<wp:wrapNone/>'
            f'<wp:docPr id="{sid}" name="ShapeKey"/>'
            f'<wp:cNvGraphicFramePr/>'
            f'<a:graphic xmlns:a="{A}">'
            f'<a:graphicData uri="{PIC_NS}">'
            f'<pic:pic xmlns:pic="{PIC_NS}">'
            f'<pic:nvPicPr>'
            f'<pic:cNvPr id="{sid}" name="ShapeKey"/>'
            f'<pic:cNvPicPr/>'
            f'</pic:nvPicPr>'
            f'<pic:blipFill>'
            f'<a:blip xmlns:r="{R_NS}" r:embed="{r_id}"/>'
            f'<a:stretch><a:fillRect/></a:stretch>'
            f'</pic:blipFill>'
            f'<pic:spPr>'
            f'<a:xfrm><a:off x="0" y="0"/><a:ext cx="{KEY_W_E}" cy="{KEY_H_E}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'</pic:spPr>'
            f'</pic:pic>'
            f'</a:graphicData>'
            f'</a:graphic>'
            f'</wp:anchor>'
            f'</w:drawing>'
        )
        run._r.append(etree.fromstring(pic_xml.encode('utf-8')))
    except Exception:
        pass

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _make_vdx(png_bytes: bytes, process_name: str, bpmn_xml: str = '') -> bytes:
    """Generate a .vdx (Visio Drawing XML) with native, fully-editable shapes.

    Parses the BPMN XML to extract every shape's position, size, type and label,
    then emits proper VDX 2D shapes (rectangles, diamonds, ellipses) and 1D
    connector arrows with arrowheads and Yes/No labels.  Falls back to an
    embedded-image VDX if the XML cannot be parsed.
    """
    import xml.etree.ElementTree as ET
    esc = _esc

    BPMN_NS   = 'http://www.omg.org/spec/BPMN/20100524/MODEL'
    BPMNDI_NS = 'http://www.omg.org/spec/BPMN/20100524/DI'
    DC_NS     = 'http://www.omg.org/spec/DD/20100524/DC'
    DI_NS     = 'http://www.omg.org/spec/DD/20100524/DI'

    # ── Fallback: embedded BMP image ─────────────────────────────────────────
    def _img_vdx() -> bytes:
        from PIL import Image as _Image
        img = _Image.open(io.BytesIO(png_bytes)).convert("RGB")
        px_w, px_h = img.size
        buf = io.BytesIO(); img.save(buf, "BMP")
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
        FDPI = 192.0; M = 0.25
        iw = px_w / FDPI; ih = px_h / FDPI
        pw = iw + 2*M;    ph = ih + 2*M
        return (
            '<?xml version="1.0" encoding="utf-8" standalone="yes"?>'
            '<VisioDocument xmlns="urn:schemas-microsoft-com:office:visio" xml:space="preserve">'
            '<StyleSheets><StyleSheet ID="0" NameU="Normal" Name="Normal"/></StyleSheets>'
            f'<DocumentSheet NameU="TheDoc"><PageProps>'
            f'<PageWidth>{pw:.4f}</PageWidth><PageHeight>{ph:.4f}</PageHeight>'
            '<PageScale>1</PageScale><DrawingScale>1</DrawingScale></PageProps></DocumentSheet>'
            '<Pages><Page ID="1" NameU="Page-1" Name="Page-1">'
            f'<PageSheet NameU="Page-1"><PageProps>'
            f'<PageWidth>{pw:.4f}</PageWidth><PageHeight>{ph:.4f}</PageHeight>'
            '<PageScale>1</PageScale><DrawingScale>1</DrawingScale></PageProps></PageSheet>'
            '<Shapes><Shape ID="1" Type="Foreign" LineStyle="0" FillStyle="0" TextStyle="0">'
            f'<XForm><PinX>{pw/2:.4f}</PinX><PinY>{ph/2:.4f}</PinY>'
            f'<Width>{iw:.4f}</Width><Height>{ih:.4f}</Height>'
            f'<LocPinX F="Width*0.5">{iw/2:.4f}</LocPinX>'
            f'<LocPinY F="Height*0.5">{ih/2:.4f}</LocPinY></XForm>'
            f'<Foreign><ImgOffsetX>0</ImgOffsetX><ImgOffsetY>0</ImgOffsetY>'
            f'<ImgWidth>{iw:.4f}</ImgWidth><ImgHeight>{ih:.4f}</ImgHeight></Foreign>'
            f'<ForeignData ForeignType="Bitmap">{b64}</ForeignData>'
            '</Shape></Shapes></Page></Pages></VisioDocument>'
        ).encode('utf-8')

    if not bpmn_xml:
        return _img_vdx()
    try:
        root = ET.fromstring(bpmn_xml)
    except ET.ParseError:
        return _img_vdx()

    # ── Collect element metadata (name + type) ────────────────────────────────
    node_info: dict[str, dict] = {}   # id → {name, type}
    flow_info: dict[str, dict] = {}   # id → {name, source, target}

    collab = root.find(f'{{{BPMN_NS}}}collaboration')
    if collab is not None:
        for el in collab:
            eid = el.get('id', '')
            if eid:
                node_info[eid] = {
                    'name': el.get('name', ''),
                    'type': el.tag.split('}')[1] if '}' in el.tag else el.tag,
                }

    proc = root.find(f'{{{BPMN_NS}}}process')
    if proc is not None:
        for el in proc.iter():
            eid = el.get('id', '')
            if not eid:
                continue
            tag = el.tag.split('}')[1] if '}' in el.tag else el.tag
            node_info[eid] = {'name': el.get('name', ''), 'type': tag}
            if tag == 'sequenceFlow':
                flow_info[eid] = {
                    'name':   el.get('name', ''),
                    'source': el.get('sourceRef', ''),
                    'target': el.get('targetRef', ''),
                }

    # ── Collect diagram bounds, edge waypoints and explicit label positions ──────
    shape_bounds:   dict[str, dict] = {}   # bpmnElement id → {x,y,w,h}
    edge_waypoints: dict[str, list] = {}   # sequenceFlow id → [(x,y), ...]
    label_bounds:   dict[str, dict] = {}   # bpmnElement id → explicit label {x,y,w,h}

    plane = root.find(f'.//{{{BPMNDI_NS}}}BPMNPlane')
    if plane is None:
        return _img_vdx()

    for child in plane:
        ctag  = child.tag.split('}')[1] if '}' in child.tag else child.tag
        bel   = child.get('bpmnElement', '')
        if ctag == 'BPMNShape':
            b = child.find(f'{{{DC_NS}}}Bounds')
            if b is not None:
                shape_bounds[bel] = {
                    'x': float(b.get('x', 0)), 'y': float(b.get('y', 0)),
                    'w': float(b.get('width',  0)), 'h': float(b.get('height', 0)),
                }
            # BPMNLabel carries the explicit position for gateway labels
            lbl_el = child.find(f'{{{BPMNDI_NS}}}BPMNLabel')
            if lbl_el is not None:
                lb = lbl_el.find(f'{{{DC_NS}}}Bounds')
                if lb is not None:
                    label_bounds[bel] = {
                        'x': float(lb.get('x', 0)), 'y': float(lb.get('y', 0)),
                        'w': float(lb.get('width',  0)), 'h': float(lb.get('height', 0)),
                    }
        elif ctag == 'BPMNEdge':
            wps = [(float(w.get('x', 0)), float(w.get('y', 0)))
                   for w in child.findall(f'{{{DI_NS}}}waypoint')]
            if len(wps) >= 2:
                edge_waypoints[bel] = wps
            # Capture explicit label position for Yes/No (and other) flow labels
            lbl_el = child.find(f'{{{BPMNDI_NS}}}BPMNLabel')
            if lbl_el is not None:
                lb = lbl_el.find(f'{{{DC_NS}}}Bounds')
                if lb is not None:
                    label_bounds[bel] = {
                        'x': float(lb.get('x', 0)), 'y': float(lb.get('y', 0)),
                        'w': float(lb.get('width',  0)), 'h': float(lb.get('height', 0)),
                    }

    if not shape_bounds:
        return _img_vdx()

    # ── Page geometry ─────────────────────────────────────────────────────────
    DPI    = 96.0    # BPMN coords are at screen 96 DPI
    MARGIN = 0.3     # inches padding

    min_x = min(b['x']           for b in shape_bounds.values())
    min_y = min(b['y']           for b in shape_bounds.values())
    max_x = max(b['x'] + b['w'] for b in shape_bounds.values())
    max_y = max(b['y'] + b['h'] for b in shape_bounds.values())

    page_w = (max_x - min_x) / DPI + 2 * MARGIN
    page_h = (max_y - min_y) / DPI + 2 * MARGIN

    def ix(px: float) -> float:   # BPMN px  →  VDX x inches
        return MARGIN + (px - min_x) / DPI

    def iy(py: float) -> float:   # BPMN px  →  VDX y inches (Y-axis flipped)
        return page_h - MARGIN - (py - min_y) / DPI

    # ── Colour palette (#RRGGBB hex — safest VDX colour format) ──────────────
    # Lanes use a distinctly tinted fill so dividing lines are always visible.
    FILL = {
        'startEvent':  '#4AA366',
        'endEvent':    '#DC3545',
        'task':        '#AFC1D6',
        'gateway':     '#F5CB5C',
        'lane':        '#E3ECF5',   # light blue-grey — clearly different from page white
        'participant': '#C8D8E8',   # slightly darker blue-grey for the outer pool
    }
    BORDER = {
        'startEvent':  '#2D7846',
        'endEvent':    '#A01E28',
        'task':        '#6482A0',
        'gateway':     '#B48C28',
        'lane':        '#7890AA',   # medium blue — clearly visible dividers
        'participant': '#5070A0',   # darker blue — strong outer pool border
    }

    def _ckey(typ: str) -> str:
        t = typ.lower()
        if t.startswith('start'): return 'startEvent'   # avoids 'sendTask' false match
        if t.startswith('end'):   return 'endEvent'     # avoids 'sendTask' false match
        if 'gateway'     in t:    return 'gateway'
        if t == 'lane':           return 'lane'
        if 'participant' in t:    return 'participant'
        return 'task'

    # ── Pre-assign integer VDX shape IDs (needed for <Connects>) ─────────────
    id_map: dict[str, int] = {}
    _ctr = [1]

    def vid(key: str) -> int:
        if key not in id_map:
            id_map[key] = _ctr[0]; _ctr[0] += 1
        return id_map[key]

    # Render order: participant → lane → tasks/events/gateways (back to front)
    def _z(eid: str) -> int:
        t = (node_info.get(eid) or {}).get('type', '').lower()
        if 'participant' in t: return 0
        if 'lane'        in t: return 1
        return 2

    sorted_ids = sorted(shape_bounds.keys(), key=_z)
    for eid in sorted_ids:
        vid(eid)
    # Pre-assign IDs for external gateway label text boxes
    for eid in label_bounds:
        if 'gateway' in (node_info.get(eid) or {}).get('type', '').lower():
            vid(f'lbl_{eid}')
    # Pre-assign IDs for edge label text boxes (Yes/No labels on flows)
    for fid in edge_waypoints:
        if flow_info.get(fid, {}).get('name') and fid in label_bounds:
            vid(f'elbl_{fid}')

    # ── Build VDX XML ─────────────────────────────────────────────────────────
    out: list[str] = []
    w = out.append

    def _xform(pin_x_in: float, pin_y_in: float, w_in: float, h_in: float) -> None:
        w('<XForm>')
        w(f'<PinX>{pin_x_in:.4f}</PinX><PinY>{pin_y_in:.4f}</PinY>')
        w(f'<Width>{w_in:.4f}</Width><Height>{h_in:.4f}</Height>')
        w(f'<LocPinX F="Width*0.5">{w_in/2:.4f}</LocPinX>')
        w(f'<LocPinY F="Height*0.5">{h_in/2:.4f}</LocPinY>')
        w('<Angle>0</Angle><FlipX>0</FlipX><FlipY>0</FlipY>')
        w('</XForm>')

    w('<?xml version="1.0" encoding="utf-8" standalone="yes"?>')
    w('<VisioDocument xmlns="urn:schemas-microsoft-com:office:visio" xml:space="preserve">')
    w(f'<DocumentProperties><Creator>POET</Creator><Title>{esc(process_name)}</Title></DocumentProperties>')
    w('<StyleSheets><StyleSheet ID="0" NameU="Normal" Name="Normal"/></StyleSheets>')
    w('<DocumentSheet NameU="TheDoc"><PageProps>')
    w(f'<PageWidth>{page_w:.4f}</PageWidth><PageHeight>{page_h:.4f}</PageHeight>')
    w('<PageScale>1</PageScale><DrawingScale>1</DrawingScale>')
    w('</PageProps></DocumentSheet>')
    w('<Pages><Page ID="1" NameU="Page-1" Name="Page-1">')
    w('<PageSheet NameU="Page-1"><PageProps>')
    w(f'<PageWidth>{page_w:.4f}</PageWidth><PageHeight>{page_h:.4f}</PageHeight>')
    w('<PageScale>1</PageScale><DrawingScale>1</DrawingScale>')
    w('</PageProps></PageSheet>')
    w('<Shapes>')

    # ── 2D shapes ─────────────────────────────────────────────────────────────
    for el_id in sorted_ids:
        b    = shape_bounds[el_id]
        info = node_info.get(el_id, {'name': el_id, 'type': ''})
        typ  = info.get('type', '')
        name = info.get('name', '')
        sid  = vid(el_id)
        ck   = _ckey(typ)
        fc   = FILL[ck]; lc = BORDER[ck]

        bx, by, bw, bh = b['x'], b['y'], b['w'], b['h']
        pin_xi = ix(bx + bw / 2)
        pin_yi = iy(by + bh / 2)
        w_in   = bw / DPI
        h_in   = bh / DPI
        lhalf  = w_in / 2
        hhalf  = h_in / 2
        t      = typ.lower()

        # Gateways, participants and lanes carry their label in a separate shape
        # so the wide lane rectangle never has text that can bleed over process shapes.
        suppress_text = 'gateway' in t or 'participant' in t or t == 'lane'

        w(f'<Shape ID="{sid}" Type="Shape" LineStyle="0" FillStyle="0" TextStyle="0">')
        _xform(pin_xi, pin_yi, w_in, h_in)
        w(f'<Fill><FillForegnd>{fc}</FillForegnd><FillBkgnd>{fc}</FillBkgnd>')
        w('<FillPattern>1</FillPattern></Fill>')

        if 'event' in t:
            lw = '0.02' if 'end' in t else '0.01'
            w(f'<Line><LineWeight>{lw}</LineWeight><LineColor>{lc}</LineColor></Line>')
            w('<Geom IX="0"><Ellipse IX="1">')
            w(f'<X F="Width*0.5">{lhalf:.4f}</X><Y F="Height*0.5">{hhalf:.4f}</Y>')
            w(f'<A F="Width*0.5">{lhalf:.4f}</A><B F="0">0</B>')
            w(f'<C F="Width">{w_in:.4f}</C><D F="Height*0.5">{hhalf:.4f}</D>')
            w('</Ellipse></Geom>')

        elif 'gateway' in t:
            w(f'<Line><LineWeight>0.015</LineWeight><LineColor>{lc}</LineColor></Line>')
            w('<Geom IX="0">')
            w(f'<MoveTo IX="1"><X F="Width*0.5">{lhalf:.4f}</X><Y F="0">0</Y></MoveTo>')
            w(f'<LineTo IX="2"><X F="Width">{w_in:.4f}</X><Y F="Height*0.5">{hhalf:.4f}</Y></LineTo>')
            w(f'<LineTo IX="3"><X F="Width*0.5">{lhalf:.4f}</X><Y F="Height">{h_in:.4f}</Y></LineTo>')
            w(f'<LineTo IX="4"><X F="0">0</X><Y F="Height*0.5">{hhalf:.4f}</Y></LineTo>')
            w(f'<LineTo IX="5"><X F="Width*0.5">{lhalf:.4f}</X><Y F="0">0</Y></LineTo>')
            w('</Geom>')

        elif 'participant' in t:
            # Outer pool border — thicker line
            w(f'<Line><LineWeight>0.02</LineWeight><LineColor>{lc}</LineColor></Line>')
            # Explicit closed-rectangle geometry (required for fill/border to render in VDX)
            w('<Geom IX="0">')
            w(f'<MoveTo IX="1"><X F="0">0</X><Y F="0">0</Y></MoveTo>')
            w(f'<LineTo IX="2"><X F="Width">{w_in:.4f}</X><Y F="0">0</Y></LineTo>')
            w(f'<LineTo IX="3"><X F="Width">{w_in:.4f}</X><Y F="Height">{h_in:.4f}</Y></LineTo>')
            w(f'<LineTo IX="4"><X F="0">0</X><Y F="Height">{h_in:.4f}</Y></LineTo>')
            w(f'<LineTo IX="5"><X F="0">0</X><Y F="0">0</Y></LineTo>')
            w('</Geom>')

        elif 'lane' in t:
            # Lane rectangle — visible dividing lines between swim lanes
            w(f'<Line><LineWeight>0.015</LineWeight><LineColor>{lc}</LineColor></Line>')
            w('<Geom IX="0">')
            w(f'<MoveTo IX="1"><X F="0">0</X><Y F="0">0</Y></MoveTo>')
            w(f'<LineTo IX="2"><X F="Width">{w_in:.4f}</X><Y F="0">0</Y></LineTo>')
            w(f'<LineTo IX="3"><X F="Width">{w_in:.4f}</X><Y F="Height">{h_in:.4f}</Y></LineTo>')
            w(f'<LineTo IX="4"><X F="0">0</X><Y F="Height">{h_in:.4f}</Y></LineTo>')
            w(f'<LineTo IX="5"><X F="0">0</X><Y F="0">0</Y></LineTo>')
            w('</Geom>')

        else:
            # Task rectangle
            w(f'<Line><LineWeight>0.01</LineWeight><LineColor>{lc}</LineColor></Line>')
            w('<Geom IX="0">')
            w(f'<MoveTo IX="1"><X F="0">0</X><Y F="0">0</Y></MoveTo>')
            w(f'<LineTo IX="2"><X F="Width">{w_in:.4f}</X><Y F="0">0</Y></LineTo>')
            w(f'<LineTo IX="3"><X F="Width">{w_in:.4f}</X><Y F="Height">{h_in:.4f}</Y></LineTo>')
            w(f'<LineTo IX="4"><X F="0">0</X><Y F="Height">{h_in:.4f}</Y></LineTo>')
            w(f'<LineTo IX="5"><X F="0">0</X><Y F="0">0</Y></LineTo>')
            w('</Geom>')

        if name and not suppress_text:
            # VDX Size is in INCHES: 10 pt = 10/72 ≈ 0.1389 in
            w('<Para IX="0"><HorzAlign>1</HorzAlign><VerticalAlign>1</VerticalAlign></Para>')
            w('<Char IX="0"><Size>0.1389</Size><Color>#222222</Color></Char>')
            w(f'<Text>{esc(name)}</Text>')
        w('</Shape>')

    # ── External text labels for gateways (separate shape, no border/fill) ────
    # Uses the BPMNLabel bounds that build_bpmn_xml already places top-left of
    # each diamond so the question text never overlaps the shape or Yes/No arrows.
    for el_id, lb in label_bounds.items():
        info = node_info.get(el_id, {})
        if 'gateway' not in info.get('type', '').lower():
            continue
        name = info.get('name', '')
        if not name:
            continue
        lsid = id_map.get(f'lbl_{el_id}')
        if lsid is None:
            continue
        lx, ly, lw_px, lh_px = lb['x'], lb['y'], lb['w'], lb['h']
        lpin_x = ix(lx + lw_px / 2)
        lpin_y = iy(ly + lh_px / 2)
        lw_in  = max(lw_px / DPI, 0.05)
        lh_in  = max(lh_px / DPI, 0.05)
        w(f'<Shape ID="{lsid}" Type="Shape" LineStyle="0" FillStyle="0" TextStyle="0">')
        _xform(lpin_x, lpin_y, lw_in, lh_in)
        w('<Fill><FillPattern>0</FillPattern></Fill>')   # transparent fill
        w('<Line><LinePattern>0</LinePattern></Line>')   # no border
        w('<Para IX="0"><HorzAlign>0</HorzAlign><VerticalAlign>0</VerticalAlign></Para>')
        w('<Char IX="0"><Size>0.1389</Size><Color>#222222</Color></Char>')
        w(f'<Text>{esc(name)}</Text>')
        w('</Shape>')

    # ── Connector shapes (2D polyline body + explicit triangle arrowhead) ────
    # VDX EndArrow is only honoured on XForm1D (1D) shapes, but XForm1D shapes
    # require a Geom whose local-space coords collapse to zero for certain
    # segment orientations, making the shape invisible.
    # Solution: render each flow as TWO 2D shapes (confirmed to render correctly):
    #   1. A 2D polyline for the full orthogonal route.
    #   2. A small filled-triangle at the last waypoint, rotated to match the
    #      last segment direction, to act as the arrowhead.
    AW = 0.10    # arrowhead base→tip length, inches
    AH = 0.065   # arrowhead base total width, inches

    for flow_id, wps in edge_waypoints.items():
        if len(wps) < 2:
            continue

        vdx_wps = [(ix(p[0]), iy(p[1])) for p in wps]
        ox, oy  = vdx_wps[0]

        xs   = [p[0] for p in vdx_wps];  ys = [p[1] for p in vdx_wps]
        bb_w = max(max(xs) - min(xs), 0.01)
        bb_h = max(max(ys) - min(ys), 0.01)

        csid = vid(f'conn_{flow_id}')

        # 1. Polyline body ──────────────────────────────────────────────────
        w(f'<Shape ID="{csid}" Type="Shape" LineStyle="0" FillStyle="0" TextStyle="0">')
        w('<XForm>')
        w(f'<PinX>{ox:.4f}</PinX><PinY>{oy:.4f}</PinY>')
        w(f'<Width>{bb_w:.4f}</Width><Height>{bb_h:.4f}</Height>')
        w('<LocPinX>0</LocPinX><LocPinY>0</LocPinY>')
        w('<Angle>0</Angle><FlipX>0</FlipX><FlipY>0</FlipY>')
        w('</XForm>')
        w('<Fill><FillPattern>0</FillPattern></Fill>')
        w('<Line><LinePattern>1</LinePattern>'
          '<LineWeight>0.015</LineWeight><LineColor>#505050</LineColor></Line>')
        w('<Geom IX="0">')
        w('<MoveTo IX="1"><X>0</X><Y>0</Y></MoveTo>')
        for seg_i, (px, py) in enumerate(vdx_wps[1:], start=2):
            w(f'<LineTo IX="{seg_i}"><X>{px-ox:.4f}</X><Y>{py-oy:.4f}</Y></LineTo>')
        w('</Geom>')
        w('</Shape>')

        # 2. Filled-triangle arrowhead ──────────────────────────────────────
        # LocPinX=AW places the local tip corner (AW, AH/2) at PinX/PinY
        # so the tip lands exactly on the last waypoint after rotation.
        end_x,  end_y  = vdx_wps[-1]
        prev_x, prev_y = vdx_wps[-2]
        arrow_angle    = math.atan2(end_y - prev_y, end_x - prev_x)

        arrow_sid = vid(f'arrowhead_{flow_id}')
        w(f'<Shape ID="{arrow_sid}" Type="Shape" LineStyle="0" FillStyle="0" TextStyle="0">')
        w('<XForm>')
        w(f'<PinX>{end_x:.4f}</PinX><PinY>{end_y:.4f}</PinY>')
        w(f'<Width>{AW:.4f}</Width><Height>{AH:.4f}</Height>')
        w(f'<LocPinX>{AW:.4f}</LocPinX><LocPinY>{AH/2:.4f}</LocPinY>')
        w(f'<Angle>{arrow_angle:.6f}</Angle>')
        w('<FlipX>0</FlipX><FlipY>0</FlipY>')
        w('</XForm>')
        w('<Fill><FillForegnd>#505050</FillForegnd>'
          '<FillBkgnd>#505050</FillBkgnd><FillPattern>1</FillPattern></Fill>')
        w('<Line><LinePattern>0</LinePattern></Line>')
        # Triangle: base-bottom(0,0) → base-top(0,AH) → tip(AW,AH/2) → close
        w('<Geom IX="0">')
        w('<MoveTo IX="1"><X>0</X><Y>0</Y></MoveTo>')
        w(f'<LineTo IX="2"><X>0</X><Y>{AH:.4f}</Y></LineTo>')
        w(f'<LineTo IX="3"><X>{AW:.4f}</X><Y>{AH/2:.4f}</Y></LineTo>')
        w('<LineTo IX="4"><X>0</X><Y>0</Y></LineTo>')
        w('</Geom>')
        w('</Shape>')

    # ── Edge label text boxes (Yes / No and other flow names) ────────────────
    # Rendered as transparent shapes at the BPMNLabel bounds computed by
    # build_bpmn_xml so they always sit right beside the exiting arrow segment.
    for fid, lb in label_bounds.items():
        fi = flow_info.get(fid, {})
        if not fi:          # not a flow — skip (gateway labels handled above)
            continue
        flabel = fi.get('name', '')
        if not flabel:
            continue
        lsid = id_map.get(f'elbl_{fid}')
        if lsid is None:
            continue
        lx, ly   = lb['x'], lb['y']
        lw_px    = max(lb['w'], 20)
        lh_px    = max(lb['h'], 14)
        lpin_x   = ix(lx + lw_px / 2)
        lpin_y   = iy(ly + lh_px / 2)
        lw_in    = lw_px / DPI
        lh_in    = lh_px / DPI
        w(f'<Shape ID="{lsid}" Type="Shape" LineStyle="0" FillStyle="0" TextStyle="0">')
        _xform(lpin_x, lpin_y, lw_in, lh_in)
        w('<Fill><FillPattern>0</FillPattern></Fill>')
        w('<Line><LinePattern>0</LinePattern></Line>')
        w('<Para IX="0"><HorzAlign>0</HorzAlign><VerticalAlign>1</VerticalAlign></Para>')
        w('<Char IX="0"><Size>0.1250</Size><Color>#444444</Color></Char>')
        w(f'<Text>{esc(flabel)}</Text>')
        w('</Shape>')

    # ── Shape key: Foreign image at bottom-right ──────────────────────────────
    try:
        from PIL import Image as _PILImg
        key_png     = _make_shape_key_image()
        KEY_W_IN_V  = 2.6
        KEY_H_IN_V  = 1.75
        # Convert PNG to BMP (VDX Foreign bitmap format)
        pil_key = _PILImg.open(io.BytesIO(key_png)).convert("RGB")
        bmp_buf = io.BytesIO()
        pil_key.save(bmp_buf, "BMP")
        key_b64 = base64.b64encode(bmp_buf.getvalue()).decode("ascii")
        # Position: bottom-right, just above the MARGIN
        key_pin_x = page_w - MARGIN - KEY_W_IN_V / 2
        key_pin_y = MARGIN + KEY_H_IN_V / 2
        key_sid   = _ctr[0]; _ctr[0] += 1
        w(f'<Shape ID="{key_sid}" Type="Foreign">')
        w(f'<XForm>')
        w(f'<PinX>{key_pin_x:.4f}</PinX><PinY>{key_pin_y:.4f}</PinY>')
        w(f'<Width>{KEY_W_IN_V:.4f}</Width><Height>{KEY_H_IN_V:.4f}</Height>')
        w(f'<LocPinX>{KEY_W_IN_V / 2:.4f}</LocPinX><LocPinY>{KEY_H_IN_V / 2:.4f}</LocPinY>')
        w('</XForm>')
        w(f'<Foreign>')
        w(f'<ImgOffsetX>0</ImgOffsetX><ImgOffsetY>0</ImgOffsetY>')
        w(f'<ImgWidth>{KEY_W_IN_V:.4f}</ImgWidth><ImgHeight>{KEY_H_IN_V:.4f}</ImgHeight>')
        w(f'</Foreign>')
        w(f'<ForeignData ForeignType="Bitmap">{key_b64}</ForeignData>')
        w('</Shape>')
    except Exception:
        pass

    # ── Lane label strips (separate narrow shapes, always left of process shapes) ─
    # LANE_LBL_W_IN = 120 px / 96 DPI = 1.25 inches — matches LANE_LBL_W constant
    # in build_bpmn_xml so there is guaranteed clear space to the right.
    LANE_LBL_W_IN = 120.0 / DPI
    for el_id, b in shape_bounds.items():
        info = node_info.get(el_id, {})
        if info.get('type', '').lower() != 'lane':
            continue
        lname = info.get('name', '')
        if not lname:
            continue
        # Pin at the centre of the label strip (left of the lane)
        lbl_pin_x = ix(b['x']) + LANE_LBL_W_IN / 2
        lbl_pin_y = iy(b['y'] + b['h'] / 2)
        lbl_h_in  = b['h'] / DPI
        lbl_sid   = _ctr[0]; _ctr[0] += 1
        w(f'<Shape ID="{lbl_sid}" Type="Shape" LineStyle="0" FillStyle="0" TextStyle="0">')
        w('<XForm>')
        w(f'<PinX>{lbl_pin_x:.4f}</PinX><PinY>{lbl_pin_y:.4f}</PinY>')
        w(f'<Width>{LANE_LBL_W_IN:.4f}</Width><Height>{lbl_h_in:.4f}</Height>')
        w(f'<LocPinX F="Width*0.5">{LANE_LBL_W_IN/2:.4f}</LocPinX>')
        w(f'<LocPinY F="Height*0.5">{lbl_h_in/2:.4f}</LocPinY>')
        w('<Angle>0</Angle><FlipX>0</FlipX><FlipY>0</FlipY>')
        w('</XForm>')
        w('<Fill><FillForegnd>#C8D8EC</FillForegnd><FillBkgnd>#C8D8EC</FillBkgnd>')
        w('<FillPattern>1</FillPattern></Fill>')
        w(f'<Line><LineWeight>0.01</LineWeight><LineColor>#7890AA</LineColor></Line>')
        w('<Geom IX="0">')
        w(f'<MoveTo IX="1"><X>0</X><Y>0</Y></MoveTo>')
        w(f'<LineTo IX="2"><X>{LANE_LBL_W_IN:.4f}</X><Y>0</Y></LineTo>')
        w(f'<LineTo IX="3"><X>{LANE_LBL_W_IN:.4f}</X><Y>{lbl_h_in:.4f}</Y></LineTo>')
        w(f'<LineTo IX="4"><X>0</X><Y>{lbl_h_in:.4f}</Y></LineTo>')
        w(f'<LineTo IX="5"><X>0</X><Y>0</Y></LineTo>')
        w('</Geom>')
        # Clip text to the label box so it never spills into the process area
        w('<TextBlock><TextMarginLeft>0.04</TextMarginLeft>')
        w('<TextMarginRight>0.04</TextMarginRight></TextBlock>')
        w('<Para IX="0"><HorzAlign>1</HorzAlign><VerticalAlign>1</VerticalAlign></Para>')
        w('<Char IX="0"><Size>0.1111</Size><Color>#2C3E50</Color><Style>1</Style></Char>')
        w(f'<Text>{esc(lname)}</Text>')
        w('</Shape>')

    w('</Shapes>')
    w('</Page></Pages></VisioDocument>')
    return ''.join(out).encode('utf-8')


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/")
def root():
    from fastapi.responses import RedirectResponse
    return RedirectResponse(url="/index.html")


@app.get("/health")
def health():
    return {"status": "ok", "service": "POET API"}


async def _extract_document_text(files: List[UploadFile]) -> str:
    """Read, validate, and extract text from a list of uploaded files."""
    text_parts = []
    for file in files:
        ext = Path(file.filename).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type '{ext}' in '{file.filename}'. Accepted: {', '.join(ALLOWED_EXTENSIONS)}"
            )
        file_bytes = await file.read()
        size_mb = len(file_bytes) / (1024 * 1024)
        if size_mb > MAX_FILE_SIZE_MB:
            raise HTTPException(
                status_code=413,
                detail=f"File '{file.filename}' is too large ({size_mb:.1f} MB). Maximum is {MAX_FILE_SIZE_MB} MB."
            )
        try:
            text = extract_text(file_bytes, file.filename)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Could not parse '{file.filename}': {str(e)}")
        if text.strip():
            text_parts.append(f"[Document: {file.filename}]\n{text}")
    if not text_parts:
        raise HTTPException(status_code=422, detail="No text could be extracted from the uploaded documents.")
    return "\n\n---\n\n".join(text_parts)


@app.post("/api/identify-processes")
async def identify_processes(files: List[UploadFile] = File(...)):
    """Scan uploaded documents and return a list of distinct processes found."""
    document_text = await _extract_document_text(files)
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY is not configured in backend/.env")
    try:
        client = anthropic.Anthropic(api_key=api_key)
        truncated = document_text[:8000] if len(document_text) > 8000 else document_text
        msg = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1024,
            messages=[{"role": "user", "content": f"{IDENTIFY_PROCESSES_PROMPT}\n\n---\n\n{truncated}"}]
        )
        raw = msg.content[0].text.strip()
        raw = re.sub(r"^```[a-zA-Z]*\n?", "", raw)
        raw = re.sub(r"\n?```$", "", raw).strip()
        match = re.search(r'\[.*\]', raw, re.DOTALL)
        processes = json.loads(match.group() if match else raw)
        return JSONResponse({"processes": processes})
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"Could not parse process list: {e}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Process identification failed: {str(e)}")


@app.post("/api/upload")
async def upload_document(files: List[UploadFile] = File(...), process_title: str = Form(""), bpmn_level: str = Form("2"), focus_process: str = Form(""), map_type: str = Form("")):
    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded.")

    document_text = await _extract_document_text(files)

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY is not configured in backend/.env")

    try:
        client = anthropic.Anthropic(api_key=api_key)
        structure = extract_process_structure(client, document_text, bpmn_level, focus_process=focus_process, map_type=map_type)
        if process_title.strip():
            structure["process_name"] = process_title.strip()
        bpmn_xml = build_bpmn_xml(structure)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"BPMN generation failed: {str(e)}")

    valid, err = validate_xml(bpmn_xml)
    if not valid:
        raise HTTPException(status_code=500, detail=f"Generated XML failed validation: {err}")

    filenames = ", ".join(f.filename for f in files)
    return JSONResponse({
        "status": "success",
        "filename": filenames,
        "process_name": structure.get("process_name", ""),
        "step_count": len(structure.get("steps", [])),
        "bpmn_xml": bpmn_xml
    })


IMPL_PLAN_SECTIONS = [
    ("executive_summary",      "Executive Summary",
     "Brief overview of the transformation and expected benefits."),
    ("gap_analysis",           "Gap Analysis",
     "Key differences between current state and future state — what needs to change."),
    ("implementation_phases",  "Implementation Phases",
     "Break the transition into clear phases (e.g. Phase 1: Discovery, Phase 2: Design, "
     "Phase 3: Implementation, Phase 4: Testing, Phase 5: Go-Live & Stabilisation). "
     "For each phase include:\n"
     "- **Objective:** What this phase achieves\n"
     "- **Key Activities:** Bullet list of tasks\n"
     "- **Deliverables:** Tangible outputs\n"
     "- **Timeline:** Suggested duration (weeks)\n"
     "- **Owner / Responsible Party:** Who leads this phase"),
    ("resource_requirements",  "Resource Requirements",
     "People, tools, systems, and budget considerations."),
    ("risks_mitigations",      "Risks & Mitigations",
     "Top risks with likelihood, impact, and mitigation actions."),
    ("success_metrics",        "Success Metrics & KPIs",
     "How to measure a successful transition."),
    ("change_management",      "Change Management",
     "Stakeholder communication, training needs, and adoption strategy."),
]

IMPL_PLAN_SECTIONS_MAP = {sid: (title, desc) for sid, title, desc in IMPL_PLAN_SECTIONS}

DEFAULT_IMPL_PLAN_SECTIONS = [sid for sid, _, _ in IMPL_PLAN_SECTIONS]

IMPLEMENTATION_PLAN_PROMPT_BASE = """You are a business transformation consultant. Based on the source documents and the process names provided, generate a structured implementation plan to transition from the Current State process to the Future State process.

Format the plan in markdown with ONLY the sections listed below (in the order given). Do not add any extra sections.

{sections_block}

Rules:
- Use ## for section headings, **Bold:** for sub-field labels
- Use numbered lists for sequential steps, bullet lists for non-sequential items
- Be specific and actionable — avoid generic filler
- Be concise: 4–8 bullet points or sentences per section maximum
- Do NOT use ALL CAPS
- Do NOT use markdown tables — use bullet lists or numbered lists instead
- Do NOT add a title or preamble — start directly with the first ## section heading
- All financial figures, costs, savings, and estimates MUST be expressed in US Dollars (USD, $). Never use GBP, £, EUR, or any other currency.
- Return only the implementation plan document"""


def _build_impl_plan_prompt(selected_sections: list) -> str:
    ordered = [s for s in IMPL_PLAN_SECTIONS if s[0] in selected_sections]
    blocks = []
    for sid, title, desc in ordered:
        blocks.append(f"## {title}\n{desc}")
    sections_block = "\n\n".join(blocks)
    return IMPLEMENTATION_PLAN_PROMPT_BASE.format(sections_block=sections_block)


IMPL_PARAMETERS = {
    "business_continuity":    ("Business Continuity",
        "Ensure no disruption to operations, with parallel systems and rollback options."),
    "customer_experience":    ("Customer Experience",
        "Maintain or improve speed, transparency, and satisfaction throughout the transition."),
    "phased_delivery":        ("Phased Delivery",
        "Implement in waves, prioritising high-volume, low-complexity use cases to deliver quick wins."),
    "human_in_the_loop":      ("Human-in-the-Loop",
        "Augment staff with automation/AI while retaining human oversight and escalation for complex decisions."),
    "data_governance":        ("Data & Governance",
        "Establish high-quality data foundations and strong governance (accuracy, fairness, auditability)."),
    "technology_flexibility": ("Technology Flexibility",
        "Use modular, API-driven architecture that integrates with legacy systems."),
    "regulatory_compliance":  ("Regulatory Compliance",
        "Ensure explainability, audit trails, and alignment with applicable regulations."),
    "change_management":      ("Change Management",
        "Redesign roles, train employees, and drive adoption through clear communication and incentives."),
    "financial_accountability":("Financial Accountability",
        "Track ROI with clear metrics and tie investment to measurable outcomes."),
    "performance_measurement":("Performance Measurement",
        "Define and monitor KPIs (e.g. cycle time, cost per unit, error rate, customer satisfaction)."),
    "scalability":            ("Scalability",
        "Build solutions that can expand across products, regions, and volumes over time."),
    "governance_leadership":  ("Governance & Leadership",
        "Establish clear ownership, decision rights, and strong executive sponsorship."),
}


async def _stream_implementation_plan(api_key: str, document_text: str,
                                       current_process: str, future_process: str,
                                       selected_sections: list = None,
                                       selected_parameters: list = None):
    if not selected_sections:
        selected_sections = DEFAULT_IMPL_PLAN_SECTIONS
    system_prompt = _build_impl_plan_prompt(selected_sections)
    params_block = ""
    if selected_parameters:
        param_lines = []
        for pid in selected_parameters:
            if pid in IMPL_PARAMETERS:
                label, desc = IMPL_PARAMETERS[pid]
                param_lines.append(f"- **{label}:** {desc}")
        if param_lines:
            params_block = (
                "\n\nThe plan must specifically address the following implementation parameters "
                "throughout the content — weave them into every relevant section:\n"
                + "\n".join(param_lines)
            )
    user_msg = (
        f"Generate an implementation plan to transition from the current state to the future state "
        f"for the following process:\n\n"
        f"- Current State Process: {current_process}\n"
        f"- Future State Process: {future_process}\n\n"
        f"Source documents:\n\n{document_text[:12000]}"
        f"{params_block}"
    )
    async_client = anthropic.AsyncAnthropic(api_key=api_key)
    async with async_client.messages.stream(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        system=system_prompt,
        messages=[{"role": "user", "content": user_msg}]
    ) as stream:
        async for text in stream.text_stream:
            yield text


@app.post("/api/generate-implementation-plan")
async def generate_implementation_plan(
    files: List[UploadFile] = File(...),
    current_process:    str = Form(""),
    future_process:     str = Form(""),
    sections_json:      str = Form("[]"),
    parameters_json:    str = Form("[]"),
):
    document_text = await _extract_document_text(files)
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY is not configured in backend/.env")
    selected_sections   = json.loads(sections_json)   if sections_json.strip()   else DEFAULT_IMPL_PLAN_SECTIONS
    selected_parameters = json.loads(parameters_json) if parameters_json.strip() else []
    if not selected_sections:
        selected_sections = DEFAULT_IMPL_PLAN_SECTIONS
    from fastapi.responses import StreamingResponse as _SR
    return _SR(
        _stream_implementation_plan(
            api_key, document_text, current_process, future_process,
            selected_sections, selected_parameters
        ),
        media_type="text/plain; charset=utf-8",
    )


@app.post("/api/export")
async def export_diagram(
    format:       str = Form(...),
    png_base64:   str = Form(...),
    process_name: str = Form("Process Map"),
    bpmn_xml:     str = Form(""),
):
    try:
        png_bytes = base64.b64decode(png_base64)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid PNG data — could not decode base64.")

    fmt = format.lower().strip()
    try:
        if fmt == "pptx":
            content  = _make_pptx(png_bytes, process_name, bpmn_xml)
            media    = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename = "process-map.pptx"
        elif fmt == "docx":
            content  = _make_docx(png_bytes, process_name, bpmn_xml)
            media    = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            filename = "process-map.docx"
        elif fmt == "vsdx":
            content  = _make_vdx(png_bytes, process_name, bpmn_xml)
            media    = "application/vnd.visio"
            filename = "process-map.vdx"
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported export format: {fmt}")
    except HTTPException:
        raise
    except Exception as exc:
        import traceback
        raise HTTPException(status_code=500, detail=f"Export error ({fmt}): {exc}\n{traceback.format_exc()}")

    return Response(
        content=content,
        media_type=media,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── SOP Generation ────────────────────────────────────────────────────────────

SOP_SYSTEM_PROMPT = """You are a professional business analyst and technical writer specialising in standard operating procedures (SOPs).

Generate a clear, well-structured SOP in markdown format based on the source documents provided.

Formatting rules:
- Use # for the SOP title only (H1)
- Use ## for section headings (H2)
- Use numbered lists (1. 2. 3.) for procedure steps
- Use bullet lists (- ) for non-sequential items
- Use **Bold:** format for key-value pairs (e.g. **Department:** Finance)
- CRITICAL: Each **Bold label:** sub-field within a step MUST start on its own NEW LINE. NEVER write two bold labels on the same line. Every time you write a **Bold label:** it must be preceded by a newline. Example of CORRECT format:
  1. **Trigger:** Description of trigger.
     **Responsible Party:** Name or role.
     **Systems/Tools Required:** List of tools.
     **Step Description:** Full description here.
     **Key Control:** Control point here.
  Example of WRONG format (DO NOT DO THIS):
  1. **Trigger:** text. **Responsible Party:** text. **Systems/Tools Required:** text.
- Do NOT use markdown tables
- NEVER write text in ALL CAPS. Every single word in the document must be in normal mixed case (e.g. "Trigger:", not "TRIGGER:"). This applies to labels, headings, content, and every other part of the document without exception. Use **bold** for emphasis only.
- All financial figures, costs, and estimates MUST be expressed in US Dollars (USD, $). Never use GBP, £, EUR, or any other currency.
- ASSUMPTION REFERENCES: Every quantitative figure in the document (cycle times, durations, dollar amounts, FTE counts, error rates, frequencies, etc.) MUST be followed immediately by an inline superscript reference marker in the format <sup>[A1]</sup>, <sup>[A2]</sup>, <sup>[A3]</sup> etc. These markers must correspond exactly to numbered **Assumption [A1]:** entries in the Sources section. Number assumptions sequentially across the entire document.
- Do NOT use strikethrough text (~~text~~). Never cross out, mark out, or use strikethrough formatting anywhere in the document.
- Do NOT use ALL CAPS text anywhere in the document. Use normal sentence case or title case only.
- Do NOT include any explanation or preamble outside the SOP itself — return only the SOP document"""

# Ordered list of all supported SOP sections
SOP_SECTIONS_META = [
    ("doc_control",             "Document Control & Governance",
     "Include: SOP Title, Unique Document ID (e.g. SOP-001), Version Number, Effective Date, Review Date, "
     "Owner (Business Function), Approver(s), Superseded Versions, Change Log with summary of revisions, "
     "Distribution List, and Classification (e.g. Confidential/Internal)."),
    ("purpose",                 "Purpose & Objective",
     "Clearly state why the SOP exists, what risk or business requirement it addresses, "
     "and any regulatory or policy drivers."),
    ("scope",                   "Scope",
     "Define boundaries precisely: business units covered, products/services covered, "
     "geographies/jurisdictions, in-scope vs out-of-scope activities, and applicable legal entities."),
    ("regulatory_refs",         "Regulatory & Policy References",
     "List applicable laws (e.g. AMLD, SEC, FINRA, FCA), internal policies (Risk, Compliance, Data Privacy), "
     "industry standards (ISO, SOC 2, PCI-DSS), and related SOPs."),
    ("definitions",             "Definitions & Acronyms",
     "Define technical terms, regulatory definitions, risk categories, and system names. "
     "Use **Term:** definition format for each entry."),
    ("roles_responsibilities",  "Roles & Responsibilities (RACI)",
     "Define accountability for: Process Owner, Business Operators, Compliance, Risk, "
     "Internal Audit, and IT/System Support. Use **Role:** responsibility format for each."),
    ("process_overview",        "Process Overview",
     "High-level summary before detailed steps: describe inputs and outputs, trigger events, "
     "key decision points, and critical control checkpoints."),
    ("procedures",              "Detailed Procedures",
     "For each step include: trigger, responsible party, required systems/tools, step description, "
     "control/check required, evidence/documentation required, SLA/timeline, and escalation trigger. "
     "Separate operational steps from control steps and clearly mark Key Controls. Use numbered list format."),
    ("controls_risk",           "Controls & Risk Management",
     "Include: key risks addressed, preventive controls, detective controls, manual vs automated controls, "
     "control frequency, and control evidence retention requirements. "
     "Use **Risk:** control format for each pair."),
    ("exception_handling",      "Exception Handling & Escalation",
     "Define: what qualifies as an exception, approval thresholds, escalation chain, "
     "regulatory breach reporting procedures, and documentation requirements."),
    ("systems_data",            "Systems & Data Requirements",
     "List systems used, data inputs, data validation checks, access control requirements, "
     "data privacy/security requirements, retention requirements, "
     "segregation of duties, and access provisioning controls."),
    ("documentation_retention", "Documentation & Record Retention",
     "Specify: what documents must be stored, where (system/location), "
     "retention period, and audit trail requirements."),
    ("kpis_monitoring",         "KPIs / Monitoring & Reporting",
     "Define: process SLAs, error rates, control failures, regulatory reporting metrics, "
     "quality assurance reviews, and dashboard ownership."),
    ("training",                "Training Requirements",
     "List: required certifications, mandatory annual training, system training, "
     "and evidence of training completion requirements."),
    ("business_continuity",     "Business Continuity & Contingency",
     "Include: backup procedures, manual fallback processes, disaster recovery steps, "
     "and communication protocols during disruptions."),
    ("appendices",              "Appendices",
     "List any relevant appendices: process flowcharts, templates, forms, checklists, "
     "sample reports, and control testing scripts."),
    ("sources",                 "Sources and Assumptions",
     "This section must cover two things:\n"
     "1. Source documents — list all documents, references, policies, regulations, and data files used to produce this SOP. "
     "For each source include the document name and, where identifiable, its version, date, or owning team. "
     "Also include any external sources (journals, websites, frameworks, standards) drawn upon. "
     "Use **Source:** description format for each entry.\n"
     "2. Metric assumptions — for every quantitative figure cited in this document (cycle times, durations, "
     "frequencies, dollar amounts, FTE counts, error rates, etc.), list a corresponding numbered assumption entry "
     "using the EXACT marker that appears inline in the document body (e.g. **Assumption [A1]:** description, "
     "**Assumption [A2]:** description). For each entry state: what the figure is, where it came from, "
     "the basis or benchmark used, and any caveats. Every inline [Ax] marker in the document MUST have a "
     "matching entry here. Do not omit any."),
    ("glossary",               "Acronym & Term Dictionary",
     "Provide a comprehensive dictionary of all acronyms, abbreviations, and industry-specific or technical terms "
     "used in this document that may be unfamiliar to a general business reader. "
     "For each entry use the format: **TERM / ACRONYM:** Full expansion and plain-English definition. "
     "Order entries alphabetically. Include at minimum: all acronyms used in the document, "
     "regulatory and compliance terms, system/tool names, and any domain-specific jargon."),
]

# Map from id → (title, instructions) for quick lookup
SOP_SECTIONS_MAP = {sid: (title, instr) for sid, title, instr in SOP_SECTIONS_META}

DEFAULT_SECTIONS = ['purpose', 'scope', 'definitions', 'roles_responsibilities',
                    'procedures', 'exception_handling', 'documentation_retention']


def _build_sop_messages(combined_text: str, sop_title: str,
                        selected_sections: list) -> tuple:
    today = date.today().strftime("%B %Y")
    ordered = [item for item in SOP_SECTIONS_META if item[0] in selected_sections]
    section_lines = []
    for i, (sid, title, instr) in enumerate(ordered, 1):
        section_lines.append(f"{i}. ## {title}\n   {instr}")
    sections_block = "\n\n".join(section_lines)
    system = (
        SOP_SYSTEM_PROMPT
        + f"\n\nThe SOP must start with:\n# [SOP Title]\n"
          f"**Version:** 1.0  **Date:** {today}\n---\n\n"
          f"Then include exactly these sections in this order:\n\n{sections_block}"
    )
    user_msg = (
        f'Generate a professional SOP titled "{sop_title}" '
        f"based on the following source documents:\n\n{combined_text[:12000]}"
    )
    return system, user_msg


async def _stream_sop(api_key: str, combined_text: str,
                      sop_title: str, selected_sections: list,
                      style_hint: str = ""):
    system, user_msg = _build_sop_messages(combined_text, sop_title, selected_sections)
    if style_hint.strip():
        user_msg += f"\n\nStyle instruction: {style_hint.strip()}"
    async_client = anthropic.AsyncAnthropic(api_key=api_key)
    async with async_client.messages.stream(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        system=system,
        messages=[{"role": "user", "content": user_msg}]
    ) as stream:
        async for text in stream.text_stream:
            yield text


def _add_runs(paragraph, text: str):
    """Add a paragraph's text, converting **bold** and <sup> spans to styled runs.

    Handles <sup> nested inside **bold** (e.g. **Phase 1<sup>[A1]</sup>:**).
    Also strips any remaining stray HTML tags.
    """
    # Strip any HTML other than <sup> (e.g. <br>, <em>) — keep <sup> for superscript
    text = re.sub(r'<(?!/?sup\b)[^>]+>', '', text)

    token_re = re.compile(r'\*\*(.+?)\*\*|<sup>(.*?)</sup>', re.DOTALL)
    cursor = 0
    for m in token_re.finditer(text):
        if m.start() > cursor:
            paragraph.add_run(text[cursor:m.start()])
        if m.group(1) is not None:          # **bold** — may contain nested <sup>
            bold_text = m.group(1)
            sup_re = re.compile(r'<sup>(.*?)</sup>', re.DOTALL)
            sc = 0
            for sm in sup_re.finditer(bold_text):
                if sm.start() > sc:
                    r = paragraph.add_run(bold_text[sc:sm.start()])
                    r.bold = True
                r = paragraph.add_run(sm.group(1))
                r.bold = True
                r.font.superscript = True
                sc = sm.end()
            if sc < len(bold_text):
                r = paragraph.add_run(bold_text[sc:])
                r.bold = True
        else:                               # <sup>...</sup>
            run = paragraph.add_run(m.group(2))
            run.font.superscript = True
        cursor = m.end()
    if cursor < len(text):
        paragraph.add_run(text[cursor:])


def _set_tab_stop(paragraph, position):
    """Add a left-aligned tab stop at `position` (an Inches/Pt/Emu value) to a paragraph."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    # position is in EMU (python-docx Length); Word XML tab pos is in twentieths-of-a-point (twips)
    twips = int(position / 914400 * 1440)
    pPr = paragraph._p.get_or_add_pPr()
    tabs = OxmlElement('w:tabs')
    tab = OxmlElement('w:tab')
    tab.set(qn('w:val'), 'left')
    tab.set(qn('w:pos'), str(twips))
    tabs.append(tab)
    pPr.append(tabs)


def _fix_caps_line(line: str) -> str:
    """If a line is predominantly uppercase, convert to mixed case preserving **bold** markers."""
    # Measure uppercase ratio of plain text (excluding markdown markers)
    plain = re.sub(r'\*\*[^*]+\*\*', '', line)
    alpha = [c for c in plain if c.isalpha()]
    if not alpha or (sum(c.isupper() for c in alpha) / len(alpha)) < 0.65:
        return line  # Already mixed case — leave untouched

    # Split the line into bold-marker segments and plain segments, fix each
    parts = re.split(r'(\*\*[^*]+\*\*)', line)
    fixed = []
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            inner = part[2:-2]
            # Convert label inside bold to Title Case (e.g. TRIGGER: → Trigger:)
            fixed.append(f'**{inner.title()}**')
        else:
            # Convert plain text to sentence case (lowercase, first char of sentence capitalised)
            fixed.append(part.lower())
    result = ''.join(fixed)
    # Re-capitalise the very first letter of the line
    for i, ch in enumerate(result):
        if ch.isalpha():
            return result[:i] + result[i].upper() + result[i+1:]
    return result


def _preprocess_sop_markdown(md: str) -> str:
    """Split inline **Bold label:** fields onto their own lines and remove ALL CAPS text.

    Only applies the bold-label line-break to plain paragraph lines — list items
    (lines starting with -, *, digits) are left untouched so bullet+bold combos
    like '- **Label:** text' are not broken into a bare '-' plus orphaned text.
    """
    lines = md.split('\n')
    out = []
    for line in lines:
        stripped = line.strip()
        # Apply inline bold-label splitting only on plain paragraph lines
        if stripped and not re.match(r'^[#\-\*\d>]', stripped):
            line = re.sub(r'([^\n])\s+(\*\*[A-Za-z][^*\n]{1,50}:\*\*)', r'\1\n   \2', line)
            line = re.sub(r'([^\n])\s+(\*\*[A-Za-z][^*\n]{1,50}:\*\*)', r'\1\n   \2', line)
        out.append(_fix_caps_line(line))
    return _move_legend_to_end(_strip_blockquotes('\n'.join(out)))


def _move_legend_to_end(md: str) -> str:
    """Move a '## Legend' section to the very end of the document."""
    lines = md.split('\n')
    legend_start = None
    legend_end   = None
    for i, ln in enumerate(lines):
        s = ln.strip()
        if re.match(r'^##\s+Legend\b', s, re.IGNORECASE) and legend_start is None:
            legend_start = i
        elif legend_start is not None and i > legend_start and re.match(r'^##\s+', s):
            legend_end = i
            break
    if legend_start is None:
        return md
    if legend_end is None:
        legend_end = len(lines)
    legend_block = lines[legend_start:legend_end]
    rest = lines[:legend_start] + lines[legend_end:]
    while rest and not rest[-1].strip():
        rest.pop()
    return '\n'.join(rest + ['', ''] + legend_block)


def _strip_blockquotes(md: str) -> str:
    """Remove leading '>' blockquote markers from lines (LLM sometimes adds them)."""
    return '\n'.join(
        re.sub(r'^>\s?', '', ln) if ln.strip().startswith('>') else ln
        for ln in md.split('\n')
    )


def _parse_md_table(lines: list) -> tuple:
    """Parse a list of markdown pipe-table lines → (headers: list[str], rows: list[list[str]])."""
    headers, rows = [], []
    for line in lines:
        cells = [c.strip() for c in line.strip().strip('|').split('|')]
        # Separator row — skip
        if all(re.match(r'^[-: ]+$', c) for c in cells if c):
            continue
        if not headers:
            headers = cells
        else:
            rows.append(cells)
    return headers, rows


def _is_table_line(s: str) -> bool:
    return s.startswith('|') and '|' in s[1:]


def _make_sop_docx(sop_markdown: str, sop_title: str) -> bytes:
    from docx import Document as DocxDocument
    from docx.shared import Pt, RGBColor, Inches

    sop_markdown = _preprocess_sop_markdown(sop_markdown)

    doc = DocxDocument()

    # ── Page setup — 1-inch margins, Letter ──────────────────────────────────
    for sec in doc.sections:
        sec.page_width   = int(8.5 * 914400)
        sec.page_height  = int(11.0 * 914400)
        sec.left_margin  = sec.right_margin  = int(1.0 * 914400)
        sec.top_margin   = sec.bottom_margin = int(1.0 * 914400)

    # ── Base (Normal) style ───────────────────────────────────────────────────
    normal = doc.styles['Normal']
    normal.font.name = 'Calibri'
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after      = Pt(6)
    normal.paragraph_format.line_spacing     = Pt(14)

    # ── Heading styles — executive palette ───────────────────────────────────
    h1 = doc.styles['Heading 1']
    h1.font.name  = 'Calibri'
    h1.font.size  = Pt(18)
    h1.font.bold  = True
    h1.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
    h1.paragraph_format.space_before    = Pt(20)
    h1.paragraph_format.space_after     = Pt(6)
    h1.paragraph_format.keep_with_next  = True

    h2 = doc.styles['Heading 2']
    h2.font.name  = 'Calibri'
    h2.font.size  = Pt(14)
    h2.font.bold  = True
    h2.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    h2.paragraph_format.space_before    = Pt(16)
    h2.paragraph_format.space_after     = Pt(4)
    h2.paragraph_format.keep_with_next  = True

    try:
        h3 = doc.styles['Heading 3']
        h3.font.name  = 'Calibri'
        h3.font.size  = Pt(12)
        h3.font.bold  = True
        h3.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
        h3.paragraph_format.space_before   = Pt(12)
        h3.paragraph_format.space_after    = Pt(3)
        h3.paragraph_format.keep_with_next = True
    except KeyError:
        pass

    # ── List styles — reset indentation so all text aligns to left margin ────
    for style_name in ['List Number', 'List Bullet', 'List Continue',
                        'List Number 2', 'List Bullet 2']:
        try:
            ls = doc.styles[style_name]
            ls.font.name = 'Calibri'
            ls.font.size = Pt(11)
            ls.paragraph_format.space_after   = Pt(3)
            ls.paragraph_format.line_spacing  = Pt(14)
            ls.paragraph_format.left_indent   = Inches(0)
            ls.paragraph_format.first_line_indent = Inches(0)
        except KeyError:
            pass

    # ── Title page ────────────────────────────────────────────────────────────
    from datetime import date as _date
    from docx.oxml.ns import qn as _qn
    from docx.oxml import OxmlElement as _OxmlElement

    # Dark navy block spanning full text width
    title_para = doc.add_paragraph()
    title_para.alignment = 1  # CENTER
    pPr = title_para._p.get_or_add_pPr()
    pBdr = _OxmlElement('w:pBdr')
    for side in ('top', 'left', 'bottom', 'right'):
        bd = _OxmlElement(f'w:{side}')
        bd.set(_qn('w:val'), 'none')
        pBdr.append(bd)
    pPr.append(pBdr)
    title_para.paragraph_format.space_before = Pt(72)
    title_para.paragraph_format.space_after  = Pt(6)
    run = title_para.add_run(sop_title or 'Document')
    run.font.name  = 'Calibri'
    run.font.size  = Pt(28)
    run.font.bold  = True
    run.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    # Thin blue divider
    div = doc.add_paragraph()
    div.paragraph_format.space_before = Pt(4)
    div.paragraph_format.space_after  = Pt(4)
    div_pPr = div._p.get_or_add_pPr()
    div_bdr = _OxmlElement('w:pBdr')
    btm = _OxmlElement('w:bottom')
    btm.set(_qn('w:val'), 'single')
    btm.set(_qn('w:sz'), '4')
    btm.set(_qn('w:space'), '1')
    btm.set(_qn('w:color'), '1E40AF')
    div_bdr.append(btm)
    div_pPr.append(div_bdr)

    # Date line
    date_para = doc.add_paragraph()
    date_para.alignment = 1  # CENTER
    date_para.paragraph_format.space_before = Pt(8)
    date_para.paragraph_format.space_after  = Pt(4)
    drun = date_para.add_run(_date.today().strftime('%B %d, %Y'))
    drun.font.name  = 'Calibri'
    drun.font.size  = Pt(12)
    drun.font.color.rgb = RGBColor(0x47, 0x55, 0x69)

    doc.add_page_break()

    # ── Helper: render a markdown table block as a Word table ─────────────────
    def _flush_word_table(tbl_lines):
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.shared import RGBColor as DocxRGB

        headers, rows = _parse_md_table(tbl_lines)
        if not headers:
            return
        n_cols = len(headers)
        n_rows = 1 + len(rows)

        tbl = doc.add_table(rows=n_rows, cols=n_cols)
        tbl.style = 'Table Grid'
        tbl.alignment = WD_TABLE_ALIGNMENT.LEFT
        tbl.autofit = False  # we set widths manually

        # Category-based column widths (hardcoded to prevent drift)
        _page_w_in  = 6.5   # usable page width in inches (8.5" letter − 2×1" margins)
        def _wt(h):
            hlen = len(re.sub(r'\*\*(.+?)\*\*', r'\1', h).strip())
            if hlen <= 3:   return 1
            if hlen <= 8:   return 3
            return 6
        if n_cols == 2:
            _cws = [_page_w_in * 0.28, _page_w_in * 0.72]
        elif n_cols <= 5:
            _wts = [_wt(h) for h in headers]
            _tot = sum(_wts)
            _cws = [_page_w_in * w / _tot for w in _wts]
        else:
            _first = _page_w_in * 0.22
            _other = (_page_w_in - _first) / max(n_cols - 1, 1)
            _cws   = [_first] + [_other] * (n_cols - 1)
        for j, col_w in enumerate(_cws):
            for cell in tbl.columns[j].cells:
                cell.width = Inches(col_w)

        # Header row — dark navy background, white bold text
        hdr_row = tbl.rows[0]
        for j, h in enumerate(headers[:n_cols]):
            cell = hdr_row.cells[j]
            cell.text = ''
            para = cell.paragraphs[0]
            run = para.add_run(re.sub(r'\*\*(.+?)\*\*', r'\1', h))
            run.bold = True
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            # Dark navy cell shading
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement('w:shd')
            shd.set(qn('w:val'), 'clear')
            shd.set(qn('w:color'), 'auto')
            shd.set(qn('w:fill'), '0F172A')
            tcPr.append(shd)

        # Data rows
        for i, row_cells in enumerate(rows):
            tr = tbl.rows[i + 1]
            is_section = (
                len(row_cells) > 1 and
                all(c == '' for c in row_cells[1:]) and
                row_cells[0].startswith('**')
            )
            fill_hex = 'E2E8F0' if is_section else ('F8FAFC' if i % 2 == 0 else 'FFFFFF')
            for j in range(n_cols):
                cell = tr.cells[j]
                cell.text = ''
                txt = row_cells[j] if j < len(row_cells) else ''
                txt = re.sub(r'\*\*(.+?)\*\*', r'\1', txt)
                para = cell.paragraphs[0]
                run = para.add_run(txt)
                run.font.size = Pt(8)
                run.bold = is_section
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                shd = OxmlElement('w:shd')
                shd.set(qn('w:val'), 'clear')
                shd.set(qn('w:color'), 'auto')
                shd.set(qn('w:fill'), fill_hex)
                tcPr.append(shd)

        doc.add_paragraph()  # spacing after table

    # ── Parse and render each line ────────────────────────────────────────────
    pending_table = []
    _skipped_first_h1 = False

    def flush_pending_table():
        if pending_table:
            _flush_word_table(pending_table)
            pending_table.clear()

    for line in sop_markdown.split('\n'):
        s = line.strip()

        # Accumulate table lines
        if _is_table_line(s):
            pending_table.append(s)
            continue

        # Flush any accumulated table before processing non-table line
        flush_pending_table()

        # Skip blank lines and lone list-marker artefacts
        if not s or s in ('-', '*', '–', '•'):
            continue

        if s == '---':
            doc.add_paragraph()
            continue

        # H1: # Title — skip first occurrence (duplicates title page)
        if s.startswith('# ') and not s.startswith('## '):
            if not _skipped_first_h1:
                _skipped_first_h1 = True
            else:
                doc.add_heading(s[2:].strip(), level=1)

        # H2: ## Section
        elif s.startswith('## ') and not s.startswith('### '):
            doc.add_heading(s[3:].strip(), level=2)

        # H3: ### Sub-section
        elif s.startswith('### '):
            doc.add_heading(s[4:].strip(), level=3)

        elif s.startswith('###'):
            text = s.lstrip('#').strip()
            if text:
                doc.add_heading(text, level=3)

        # Numbered list — bold blue number, tab-stop aligned hanging indent
        elif re.match(r'^\d+\.\s', s):
            num_match = re.match(r'^(\d+)\.\s+(.*)', s)
            num_label = (num_match.group(1) + '.\t') if num_match else '1.\t'
            text      = num_match.group(2) if num_match else re.sub(r'^\d+\.\s+', '', s)
            p = doc.add_paragraph()
            pf = p.paragraph_format
            pf.left_indent        = Inches(0.28)
            pf.first_line_indent  = Inches(-0.28)
            pf.space_before       = Pt(8)
            pf.space_after        = Pt(4)
            pf.line_spacing       = Pt(14)
            _set_tab_stop(p, Inches(0.28))
            nr = p.add_run(num_label)
            nr.font.name = 'Calibri'; nr.font.size = Pt(11)
            nr.font.bold = True
            nr.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
            _add_runs(p, text)

        # Bullet list — tab-stop aligned hanging indent so all wrapped lines flush
        elif s.startswith('- ') or s.startswith('* '):
            p = doc.add_paragraph()
            pf = p.paragraph_format
            pf.left_indent        = Inches(0.22)
            pf.first_line_indent  = Inches(-0.22)
            pf.space_before       = Pt(4)
            pf.space_after        = Pt(4)
            pf.line_spacing       = Pt(14)
            _set_tab_stop(p, Inches(0.22))
            br = p.add_run('•\t')
            br.font.name = 'Calibri'; br.font.size = Pt(11)
            br.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)
            _add_runs(p, s[2:])

        # Continuation sub-field line — use Normal to avoid inherited indent
        elif s.startswith('**') and ':' in s:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0)
            _add_runs(p, s)

        # Normal paragraph
        else:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0)
            _add_runs(p, s)

    flush_pending_table()  # flush any trailing table

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@app.post("/api/generate-sop")
async def generate_sop_endpoint(
    files:      List[UploadFile] = File(...),
    sop_title:     str = Form(...),
    sections_json: str = Form("[]"),
    style_hint:    str = Form(""),
):
    if not files:
        raise HTTPException(status_code=400, detail="At least one file is required.")

    all_texts = []
    for f in files:
        ext = Path(f.filename).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type '{ext}' in {f.filename}. Accepted: {', '.join(ALLOWED_EXTENSIONS)}"
            )
        file_bytes = await f.read()
        size_mb    = len(file_bytes) / (1024 * 1024)
        if size_mb > MAX_FILE_SIZE_MB:
            raise HTTPException(
                status_code=413,
                detail=f"File '{f.filename}' is too large ({size_mb:.1f} MB). Maximum is {MAX_FILE_SIZE_MB} MB."
            )
        try:
            text = extract_text(file_bytes, f.filename)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Could not parse '{f.filename}': {str(e)}")
        if text.strip():
            all_texts.append(f"=== Source: {f.filename} ===\n{text.strip()}")

    if not all_texts:
        raise HTTPException(status_code=422, detail="No text could be extracted from the uploaded documents.")

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY is not configured in backend/.env")

    combined          = "\n\n".join(all_texts)
    selected_sections = json.loads(sections_json) if sections_json.strip() else DEFAULT_SECTIONS
    if not selected_sections:
        selected_sections = DEFAULT_SECTIONS
    if 'sources' not in selected_sections:
        selected_sections = selected_sections + ['sources']

    from fastapi.responses import StreamingResponse as _SR
    return _SR(
        _stream_sop(api_key, combined, sop_title, selected_sections, style_hint),
        media_type="text/plain; charset=utf-8",
    )


@app.post("/api/export-sop")
async def export_sop(
    sop_markdown: str = Form(...),
    sop_title:    str = Form("SOP"),
):
    try:
        content = _make_sop_docx(sop_markdown, sop_title)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Word export failed: {str(e)}")

    safe_name = re.sub(r'[^a-z0-9_\-]', '_', sop_title.lower()) or 'sop'
    return Response(
        content=content,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}.docx"'},
    )


# ── Business Case Generation ──────────────────────────────────────────────────

BC_SECTIONS_META = [
    ("executive_summary",       "Executive Summary",
     "A concise standalone overview (2-3 paragraphs) for senior stakeholders covering: the current process, "
     "proposed improvements, expected benefits, investment required, and recommendation."),
    ("current_state",           "Current State",
     "Describe the current process in detail: manual steps, pain points, inefficiencies, error rates, "
     "cycle times, FTE effort, and cost of poor quality. Reference specifics from the uploaded documents."),
    ("target_state",            "Target State",
     "Describe the proposed future-state process after implementing improvements. Cover: redesigned process flow, "
     "automation touchpoints, changes to roles, and key differences from current state."),
    ("risk_control",            "Risk & Control Impact",
     "Analyse how the proposed improvements affect existing controls. Identify controls that will be "
     "automated or enhanced, any new risks introduced, and residual manual controls required."),
    ("regulatory",              "Regulatory Considerations",
     "Identify regulatory, compliance, or legal implications. Reference applicable regulations, "
     "internal policies, and any approval or notification requirements before implementation."),
    ("financial_impact",        "Financial Impact",
     "Provide a quantified ROI analysis including: estimated annual cost savings, FTE reduction or "
     "redeployment, one-off implementation costs, ongoing operational savings, and payback period. "
     "Use **Label:** value format for key figures."),
    ("implementation_plan",     "Implementation Plan",
     "Outline a phased implementation approach with: key phases, milestones, indicative timeline, "
     "resource requirements, dependencies, and change management considerations. Use numbered list format."),
    ("technology_requirements", "Technology Requirements",
     "List the systems, tools, platforms, and integrations required. Include infrastructure changes, "
     "licensing considerations, vendor options, and build vs buy assessment."),
    ("governance",              "Governance & Ownership",
     "Define the governance structure: process owner, project sponsor, steering committee, "
     "and ongoing operational ownership post-implementation."),
    ("risks_mitigations",       "Risks & Mitigations",
     "Identify the top risks to the initiative with mitigations. Cover: delivery risks, adoption risks, "
     "regulatory risks, and operational risks. Use **Risk:** mitigation format for each pair."),
    ("sources",                 "Sources and Assumptions",
     "This section must cover two things:\n"
     "1. Source documents — list ALL sources used to produce this business case, including: uploaded documents "
     "(note each document name, type, and key information contributed) and any external sources such as industry "
     "research, academic journals, analyst reports, frameworks, standards, regulations, or websites drawn upon to "
     "support claims, benchmarks, or recommendations (include title, author/organisation, and URL or publication "
     "reference where available). Use **Source:** description format for each entry.\n"
     "2. Metric assumptions — for every quantitative figure cited in this document (cost savings, ROI, dollar "
     "amounts, FTE counts, cycle times, processing times, error rates, payback periods, percentages, etc.), list "
     "a corresponding numbered assumption entry using the EXACT marker that appears inline in the document body "
     "(e.g. **Assumption [A1]:** description, **Assumption [A2]:** description). For each entry state: what the "
     "figure is, where it came from, the basis or benchmark used, and any caveats or limitations. Every inline "
     "[Ax] marker in the document MUST have a matching entry here. Do not omit any metric."),
    ("glossary",               "Acronym & Term Dictionary",
     "Provide a comprehensive dictionary of all acronyms, abbreviations, and industry-specific or technical terms "
     "used in this document that may be unfamiliar to a general business reader. "
     "For each entry use the format: **TERM / ACRONYM:** Full expansion and plain-English definition. "
     "Order entries alphabetically. Include at minimum: all acronyms used in the document, "
     "regulatory and compliance terms, system/tool names, and any domain-specific jargon."),
]

BC_SECTIONS_MAP = {sid: (title, instr) for sid, title, instr in BC_SECTIONS_META}

BC_DEFAULT_SECTIONS = ['executive_summary', 'current_state', 'target_state',
                       'financial_impact', 'implementation_plan', 'risks_mitigations', 'sources']

IMPROVEMENT_FOCUS_META = {
    "eliminate":   "Eliminate (ESOAR): Identify and remove non-value-adding steps, redundant activities, "
                   "unnecessary approvals, and process waste. Apply value stream analysis to surface steps "
                   "that add cost or delay without delivering customer or business value.",
    "standardize": "Standardize (ESOAR): Establish consistent, repeatable process templates and controls. "
                   "Identify process variations, inconsistencies across teams or locations, and opportunities "
                   "to introduce standard operating procedures, checklists, and governance frameworks.",
    "optimize":    "Optimize (ESOAR): Improve throughput, quality, and efficiency within the existing process. "
                   "Identify bottlenecks, SLA breaches, handoff delays, rework loops, and opportunities for "
                   "LEAN continuous improvement, workload balancing, and skill-to-task alignment.",
    "automate":    "Automate (ESOAR): Apply rules-based workflow automation, business rules engines, and "
                   "system-triggered actions to reduce manual intervention. Identify decision points, "
                   "approvals, notifications, and data routing that can be handled by digital workflows.",
    "robotize":    "Robotize (ESOAR): Apply Robotic Process Automation (RPA) or AI-driven intelligent "
                   "automation to high-volume, repetitive tasks. Identify candidates for software bots, "
                   "document processing AI, NLP, predictive analytics, and cognitive automation.",
}

BC_SYSTEM_PROMPT = """You are a senior management consultant specialising in process optimisation and automation.

Generate a professional Process Optimisation Business Case in markdown format based on the source documents provided.

Formatting rules:
- Use # for the document title only (H1)
- Use ## for section headings (H2)
- Use numbered lists (1. 2. 3.) for sequential items such as steps and phases
- Use bullet lists (- ) for non-sequential items
- Use **Bold:** format for key-value pairs and financial figures
- All financial figures, costs, savings, ROI, and estimates MUST be expressed in US Dollars (USD, $). Never use GBP, £, EUR, or any other currency.
- ASSUMPTION REFERENCES: Every quantitative figure in the document (dollar amounts, cost savings, ROI, FTE counts, cycle times, processing times, error rates, payback periods, percentages, etc.) MUST be followed immediately by an inline superscript reference marker in the format <sup>[A1]</sup>, <sup>[A2]</sup>, <sup>[A3]</sup> etc. These markers must correspond exactly to numbered **Assumption [A1]:** entries in the Sources section. Number assumptions sequentially across the entire document.
- Do NOT use strikethrough text (~~text~~). Never cross out, mark out, or use strikethrough formatting anywhere in the document.
- Do NOT use ALL CAPS text anywhere in the document. Use normal sentence case or title case only.
- Do NOT use markdown tables
- Do NOT include any explanation outside the document itself — return only the business case"""


def _generate_business_case(client: anthropic.Anthropic, combined_text: str,
                             process_name: str,
                             selected_focuses: list, selected_sections: list) -> str:
    today = date.today().strftime("%B %Y")

    # Build focus instructions
    focus_lines = []
    for fid in selected_focuses:
        if fid in IMPROVEMENT_FOCUS_META:
            focus_lines.append(f"- {IMPROVEMENT_FOCUS_META[fid]}")
    focus_block = "\n".join(focus_lines)

    # Build ordered section instructions
    ordered = [item for item in BC_SECTIONS_META if item[0] in selected_sections]
    section_lines = []
    for i, (sid, title, instr) in enumerate(ordered, 1):
        section_lines.append(f"{i}. ## {title}\n   {instr}")
    sections_block = "\n\n".join(section_lines)

    system = (
        BC_SYSTEM_PROMPT
        + f"\n\nThe document must start with:\n# Process Optimisation Business Case: [Process Name]\n"
          f"**Prepared by:** POET  **Date:** {today}\n---\n\n"
          f"Improvement focus areas to analyse:\n{focus_block}\n\n"
          f"Then include exactly these sections in this order:\n\n{sections_block}"
    )

    user_msg = (
        f'Generate a business case for the "{process_name}" process. '
        f"Focus areas: {', '.join(selected_focuses)}. "
        f"Based on the following source documents:\n\n{combined_text[:12000]}"
    )

    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        system=system,
        messages=[{"role": "user", "content": user_msg}]
    )
    text = message.content[0].text
    # Strip any strikethrough markdown (~~text~~) the model may have produced
    text = re.sub(r'~~(.+?)~~', r'\1', text, flags=re.DOTALL)
    return text


def _make_bc_pptx(bc_markdown: str, process_name: str) -> bytes:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as PptxRGB

    bc_markdown = _move_legend_to_end(bc_markdown)

    prs = PptxPresentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # ── Cover slide (always generated from process_name, not from markdown H1) ─
    cover = prs.slides.add_slide(blank_layout)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = PptxRGB(0x0F, 0x17, 0x2A)
    top_bar = cover.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = PptxRGB(0x1E, 0x40, 0xAF)
    top_bar.line.fill.background()
    ttx = cover.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.6))
    ttf = ttx.text_frame; ttf.word_wrap = True
    trn = ttf.paragraphs[0].add_run()
    trn.text = process_name
    trn.font.bold = True; trn.font.size = Pt(36)
    trn.font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF)
    stx = cover.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.8))
    stf = stx.text_frame
    srn = stf.paragraphs[0].add_run()
    srn.text = date.today().strftime('%B %d, %Y')
    srn.font.size = Pt(14); srn.font.color.rgb = PptxRGB(0x93, 0xC5, 0xFD)
    bot_bar = cover.shapes.add_shape(1, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12))
    bot_bar.fill.solid(); bot_bar.fill.fore_color.rgb = PptxRGB(0x1E, 0x40, 0xAF)
    bot_bar.line.fill.background()

    lines   = bc_markdown.split('\n')
    current_title   = None
    current_content = []

    # ── PPT table layout constants ─────────────────────────────────────────────
    EMU             = 914400          # EMUs per inch
    HDR_H_EMU       = int(0.38*EMU)  # fixed header row height
    ROW_H_EMU       = int(0.28*EMU)  # fixed data row height
    TBL_TOP_IN      = 1.05           # table top y (inches)
    ROWS_PER_PAGE   = 12             # hard limit — guaranteed to fit on one slide

    def _set_row_height(tbl, row_idx, h_emu):
        """Force a fixed row height using the python-pptx public API."""
        tbl.rows[row_idx].height = h_emu

    def _make_tbl_slide(slide_title, headers, data_rows, page_num, total_pages):
        """Add one slide containing a paginated table chunk with fixed row heights."""
        from pptx.util import Pt as PPt
        sl  = prs.slides.add_slide(blank_layout)
        label = slide_title if total_pages == 1 else f"{slide_title}  ({page_num}/{total_pages})"

        # Subtle background
        bg = sl.background; bg.fill.solid()
        bg.fill.fore_color.rgb = PptxRGB(0xF8, 0xFA, 0xFC)

        # Left accent bar
        sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(7.5)).fill.solid()
        sl.shapes[-1].fill.fore_color.rgb = PptxRGB(0x1E, 0x40, 0xAF)
        sl.shapes[-1].line.fill.background()

        # Title
        tx = sl.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.73), Inches(0.7))
        tf = tx.text_frame; tf.word_wrap = False
        rn = tf.paragraphs[0].add_run()
        rn.text = label; rn.font.bold = True; rn.font.size = Pt(16)
        rn.font.color.rgb = PptxRGB(0x0F, 0x17, 0x2A)

        # Divider
        div = sl.shapes.add_shape(1, Inches(0.3), Inches(0.82), Inches(12.73), Inches(0.025))
        div.fill.solid(); div.fill.fore_color.rgb = PptxRGB(0x1E, 0x40, 0xAF)
        div.line.fill.background()

        n_cols  = len(headers)
        n_rows  = 1 + len(data_rows)
        tbl_h   = HDR_H_EMU + ROW_H_EMU * len(data_rows)
        tbl_shp = sl.shapes.add_table(n_rows, n_cols,
                                       Inches(0.3), Inches(TBL_TOP_IN),
                                       Inches(12.73), tbl_h)
        tbl = tbl_shp.table

        # Fix row heights explicitly so the table never overflows
        _set_row_height(tbl, 0, HDR_H_EMU)
        for i in range(1, n_rows):
            _set_row_height(tbl, i, ROW_H_EMU)

        # Column widths: category-based weights (hardcoded to prevent drift)
        total_emu_w = int(12.73 * EMU)
        def _col_weight(h):
            hlen = len(re.sub(r'\*\*(.+?)\*\*', r'\1', h).strip())
            if hlen <= 3:   return 1   # "#", "R", "A" etc — very narrow
            if hlen <= 8:   return 3   # short labels
            return 6                   # long labels get more space
        if n_cols == 2:
            # 2-col: narrow first col (e.g. "#" or "Ref"), wide second
            col_ws = [int(total_emu_w * 0.28), int(total_emu_w * 0.72)]
        elif n_cols <= 5:
            _wts   = [_col_weight(h) for h in headers]
            _tot   = sum(_wts)
            col_ws = [int(total_emu_w * w / _tot) for w in _wts]
        else:
            # Wide tables (RACI etc): first col slightly wider, rest equal
            first_w = int(total_emu_w * 0.22)
            other_w = (total_emu_w - first_w) // max(n_cols - 1, 1)
            col_ws  = [first_w] + [other_w] * (n_cols - 1)
        col_ws[-1] = total_emu_w - sum(col_ws[:-1])  # fix rounding
        for j in range(n_cols):
            tbl.columns[j].width = col_ws[j]

        # Header row
        for j, h in enumerate(headers[:n_cols]):
            cell = tbl.cell(0, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = PptxRGB(0x0F, 0x17, 0x2A)
            tf2 = cell.text_frame; tf2.word_wrap = True
            rn2 = tf2.paragraphs[0].add_run()
            rn2.text = re.sub(r'\*\*(.+?)\*\*', r'\1', h)
            rn2.font.bold = True; rn2.font.size = PPt(8)
            rn2.font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF)

        # Data rows
        for i, row_cells in enumerate(data_rows):
            is_sec = (len(row_cells) > 1 and
                      all(c == '' for c in row_cells[1:]) and
                      row_cells[0].startswith('**'))
            bg = (PptxRGB(0xCB, 0xD5, 0xE1) if is_sec else
                  (PptxRGB(0xF1, 0xF5, 0xF9) if i % 2 == 0 else PptxRGB(0xFF, 0xFF, 0xFF)))
            for j in range(n_cols):
                cell = tbl.cell(i + 1, j)
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                tf2 = cell.text_frame; tf2.word_wrap = True
                txt = row_cells[j] if j < len(row_cells) else ''
                txt = re.sub(r'\*\*(.+?)\*\*', r'\1', txt)
                rn2 = tf2.paragraphs[0].add_run()
                rn2.text = txt
                rn2.font.size = PPt(8); rn2.font.bold = is_sec
                rn2.font.color.rgb = PptxRGB(0x1E, 0x29, 0x3B)

    def _render_pptx_table(title, tbl_lines):
        """Parse markdown table and paginate across as many slides as needed."""
        headers, rows = _parse_md_table(tbl_lines)
        if not headers or not rows:
            return
        chunks = [rows[i:i + ROWS_PER_PAGE] for i in range(0, len(rows), ROWS_PER_PAGE)]
        total  = len(chunks)
        for pg, chunk in enumerate(chunks, start=1):
            _make_tbl_slide(title, headers, chunk, pg, total)

    def _make_text_slide(title, content_lines):
        """Create a single polished text slide. Returns the slide object."""
        from pptx.enum.text import MSO_AUTO_SIZE
        sl = prs.slides.add_slide(blank_layout)

        # Subtle light background
        bg = sl.background; bg.fill.solid()
        bg.fill.fore_color.rgb = PptxRGB(0xF8, 0xFA, 0xFC)

        # Accent bar on left edge
        sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.07), Inches(7.5)).fill.solid()
        sl.shapes[-1].fill.fore_color.rgb = PptxRGB(0x1E, 0x40, 0xAF)
        sl.shapes[-1].line.fill.background()

        # Title
        tx = sl.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.75))
        tf = tx.text_frame; tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        rn = tf.paragraphs[0].add_run()
        rn.text = re.sub(r'\*\*(.+?)\*\*', r'\1', title)
        rn.font.bold = True; rn.font.size = Pt(20)
        rn.font.color.rgb = PptxRGB(0x0F, 0x17, 0x2A)

        # Divider line
        div = sl.shapes.add_shape(1, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.025))
        div.fill.solid(); div.fill.fore_color.rgb = PptxRGB(0x1E, 0x40, 0xAF)
        div.line.fill.background()

        # Content
        cx = sl.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9))
        cf = cx.text_frame; cf.word_wrap = True
        cf.auto_size = MSO_AUTO_SIZE.NONE
        first = True
        for line in content_lines:
            s = line.strip()
            if not s:
                continue
            para = cf.paragraphs[0] if first else cf.add_paragraph()
            first = False
            is_bullet = s.startswith('- ') or s.startswith('• ') or s.startswith('* ')
            is_h3     = s.startswith('### ') or s.startswith('## ')
            if is_bullet:
                s = '• ' + re.sub(r'^[-•\*]\s+', '', s)
            elif is_h3:
                s = re.sub(r'^#+\s*', '', s)
            # Strip HTML tags — <sup>[A1]</sup> → [A1], all others removed
            s = re.sub(r'<sup>(.*?)</sup>', r'\1', s)
            s = re.sub(r'<[^>]+>', '', s)
            s = re.sub(r'\*\*(.+?)\*\*', r'\1', s)
            rn2 = para.add_run()
            rn2.text = s
            if is_h3:
                rn2.font.bold = True; rn2.font.size = Pt(13)
                rn2.font.color.rgb = PptxRGB(0x1E, 0x40, 0xAF)
                para.space_before = Pt(10)
                para.space_after  = Pt(4)
            elif is_bullet:
                rn2.font.size = Pt(10)
                rn2.font.color.rgb = PptxRGB(0x1E, 0x29, 0x3B)
                para.space_before = Pt(5)
                para.space_after  = Pt(5)
            else:
                rn2.font.size = Pt(10)
                rn2.font.color.rgb = PptxRGB(0x1E, 0x29, 0x3B)
                para.space_before = Pt(3)
                para.space_after  = Pt(3)
            para.line_spacing = 1.15
        return sl

    def flush_slide():
        if current_title is None:
            return
        # Skip if nothing to render (e.g., section header with no content, or after pagination)
        if not any(l.strip() for l in current_content):
            return

        text_lines  = [l for l in current_content if not _is_table_line(l.strip()) and l.strip()]
        table_block = [l for l in current_content if _is_table_line(l.strip())]

        if table_block:
            if text_lines:
                _make_text_slide(current_title, text_lines)
            _render_pptx_table(current_title, table_block)
        else:
            _make_text_slide(current_title, current_content)

    LINES_PER_SLIDE = 16  # max content lines per slide before paginating

    for line in lines:
        s = line.strip()
        if s.startswith('# ') and not s.startswith('## '):
            # Cover already added above from process_name — skip any H1 in markdown
            continue
        elif s.startswith('## ') or s.startswith('### '):
            flush_slide()
            current_title   = re.sub(r'^#+\s*', '', s)
            current_content = []
        elif s == '---':
            continue
        else:
            if current_title is not None:
                current_content.append(s)
                # Paginate when content grows too long
                if len([l for l in current_content if l.strip()]) >= LINES_PER_SLIDE:
                    _make_text_slide(current_title, current_content)
                    current_content = []
                    current_title   = current_title  # keep same title for continuation

    flush_slide()

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@app.post("/api/generate-business-case")
async def generate_business_case_endpoint(
    files:         List[UploadFile] = File(...),
    process_name:  str = Form(...),
    focuses_json:  str = Form("[]"),
    sections_json: str = Form("[]"),
):
    if not files:
        raise HTTPException(status_code=400, detail="At least one file is required.")

    all_texts = []
    for f in files:
        ext = Path(f.filename).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type '{ext}' in {f.filename}."
            )
        file_bytes = await f.read()
        size_mb    = len(file_bytes) / (1024 * 1024)
        if size_mb > MAX_FILE_SIZE_MB:
            raise HTTPException(
                status_code=413,
                detail=f"File '{f.filename}' is too large ({size_mb:.1f} MB)."
            )
        try:
            text = extract_text(file_bytes, f.filename)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Could not parse '{f.filename}': {str(e)}")
        if text.strip():
            all_texts.append(f"=== Source: {f.filename} ===\n{text.strip()}")

    if not all_texts:
        raise HTTPException(status_code=422, detail="No text could be extracted from the uploaded documents.")

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY is not configured in backend/.env")

    selected_focuses  = json.loads(focuses_json)  if focuses_json.strip()  else []
    selected_sections = json.loads(sections_json) if sections_json.strip() else BC_DEFAULT_SECTIONS
    if not selected_sections:
        selected_sections = BC_DEFAULT_SECTIONS
    if 'sources' not in selected_sections:
        selected_sections = selected_sections + ['sources']

    try:
        client      = anthropic.Anthropic(api_key=api_key)
        combined    = "\n\n".join(all_texts)
        bc_markdown = _generate_business_case(
            client, combined, process_name, selected_focuses, selected_sections
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Business case generation failed: {str(e)}")

    return JSONResponse({"status": "success", "bc_markdown": bc_markdown})


@app.post("/api/export-business-case")
async def export_business_case(
    bc_markdown:  str = Form(...),
    process_name: str = Form("Business Case"),
    format:       str = Form("docx"),
):
    safe_name = re.sub(r'[^a-z0-9_\-]', '_', process_name.lower()) or 'business_case'
    fmt = format.lower().strip()
    try:
        if fmt == "pptx":
            content   = _make_bc_pptx(bc_markdown, process_name)
            media     = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename  = f"{safe_name}.pptx"
        else:
            content   = _make_sop_docx(bc_markdown, process_name)
            media     = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            filename  = f"{safe_name}.docx"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

    return Response(
        content=content,
        media_type=media,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# ── Process Health Scorecard ──────────────────────────────────────────────────

SCORECARD_PROMPT = """You are a senior process excellence consultant. Analyse the uploaded process documentation and produce a Process Health Scorecard.

FORMAT RULES — follow exactly:
- Use markdown with ## for section headings, ### for sub-headings, **bold** for labels
- Do NOT use strikethrough text (~~text~~) under any circumstances
- Do NOT use ALL CAPS text anywhere in the document. Use normal sentence case or title case only.
- Do NOT include a title or preamble — start directly with the first section
- All financial figures must be in US dollars ($)
- For every metric cited, add a superscript reference marker inline: value<sup>[A1]</sup>
- The Legend section MUST appear as the very last section in the document, after Sources and Assumptions

STRUCTURE — include only the selected dimensions:
1. ## Executive Health Summary — overall RAG rating (Red / Amber / Green), 3–5 sentence narrative
2. ## Scorecard by Dimension — for each selected dimension:
   ### [Dimension Name] — Rating: [RED / AMBER / GREEN]
   - Current state observation
   - Key issues identified
   - Recommended improvement action
3. ## Priority Improvement Opportunities — top 3–5 ranked opportunities with estimated impact and priority (HIGH / MEDIUM / LOW)
4. ## Sources and Assumptions — numbered list [A1], [A2]... with derivation detail
5. ## Legend — ALWAYS include this as the very last section. Do NOT use bullet points or list markers for legend items — write each as a plain paragraph line:
**RED** — Significant problems identified. Urgent attention and remediation required. Process is materially inefficient, high-risk, or broken in this area.
**AMBER** — Notable issues present. Improvement is recommended. Process functions but has clear gaps, inefficiencies, or risks that should be addressed.
**GREEN** — Performing well. Minor improvements may be beneficial but no urgent action required.
**High / Medium / Low** (where used) — Indicates the priority or severity of an improvement opportunity: High = address immediately, Medium = plan within current cycle, Low = monitor or address when capacity allows.
All ratings are based on evidence from the uploaded source documents and consultant judgement where document data is limited.

RATING CRITERIA:
- GREEN: performing well, minor improvements only
- AMBER: notable issues, improvement recommended
- RED: significant problems, urgent attention required

Be specific and evidence-based — reference actual content from the uploaded documents."""


def _generate_scorecard(client: anthropic.Anthropic, document_text: str, process_name: str,
                        industry: str, dimensions: list, include_glossary: bool = False) -> str:
    dim_labels = {
        'waste':           'Waste & Non-Value-Adding Activity (LEAN)',
        'handoffs':        'Handoffs & Waiting Time',
        'automation':      'Automation Potential (ESOAR)',
        'rework':          'Rework & Exception Rate',
        'standardisation': 'Standardisation & Consistency',
        'controls':        'Controls & Risk Exposure',
        'data_quality':    'Data Quality & Availability',
        'customer':        'Customer / Stakeholder Impact',
    }
    dims_text = "\n".join(f"- {dim_labels.get(d, d)}" for d in dimensions)
    industry_note = f" The process operates in the {industry} sector." if industry else ""
    glossary_note = (
        "\n\nAt the end of the document, include a ## Acronym & Term Dictionary section listing "
        "all acronyms, abbreviations, and industry-specific terms used, in alphabetical order, "
        "using the format **TERM / ACRONYM:** plain-English definition."
        if include_glossary else ""
    )
    user_msg = (
        f"Process name: {process_name}.{industry_note}\n\n"
        f"Assess the following dimensions only:\n{dims_text}\n\n"
        f"Source documents:\n\n{document_text[:12000]}"
        f"{glossary_note}"
    )
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        system=SCORECARD_PROMPT,
        messages=[{"role": "user", "content": user_msg}],
    )
    return msg.content[0].text


@app.post("/api/generate-scorecard")
async def generate_scorecard_endpoint(
    files:            List[UploadFile] = File(...),
    process_name:     str = Form(...),
    industry:         str = Form(""),
    dimensions_json:  str = Form("[]"),
    include_glossary: str = Form("false"),
):
    if not files:
        raise HTTPException(status_code=400, detail="At least one file is required.")
    all_texts = []
    for f in files:
        ext = Path(f.filename).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(status_code=400, detail=f"Unsupported file type '{ext}'.")
        fb = await f.read()
        if len(fb) / 1024 / 1024 > MAX_FILE_SIZE_MB:
            raise HTTPException(status_code=413, detail=f"File '{f.filename}' too large.")
        try:
            t = extract_text(fb, f.filename)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Could not parse '{f.filename}': {e}")
        if t.strip():
            all_texts.append(f"=== {f.filename} ===\n{t.strip()}")
    if not all_texts:
        raise HTTPException(status_code=422, detail="No text could be extracted.")
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured.")
    dimensions = json.loads(dimensions_json) if dimensions_json.strip() else []
    if not dimensions:
        dimensions = ['waste', 'handoffs', 'automation', 'rework', 'standardisation', 'controls', 'data_quality', 'customer']
    try:
        client = anthropic.Anthropic(api_key=api_key)
        md = _generate_scorecard(client, "\n\n".join(all_texts), process_name, industry, dimensions,
                                 include_glossary=include_glossary.lower() == 'true')
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Scorecard generation failed: {e}")
    md = _move_legend_to_end(_strip_blockquotes(md))
    return JSONResponse({"status": "success", "sc_markdown": md})


@app.post("/api/export-scorecard")
async def export_scorecard(
    sc_markdown:  str = Form(...),
    process_name: str = Form("Process Health Scorecard"),
    format:       str = Form("docx"),
):
    safe_name = re.sub(r'[^a-z0-9_\-]', '_', process_name.lower()) or 'scorecard'
    fmt = format.lower().strip()
    try:
        if fmt == "pptx":
            content  = _make_bc_pptx(sc_markdown, process_name)
            media    = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename = f"{safe_name}.pptx"
        else:
            content  = _make_sop_docx(sc_markdown, process_name)
            media    = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            filename = f"{safe_name}.docx"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Export failed: {e}")
    return Response(content=content, media_type=media,
                    headers={"Content-Disposition": f'attachment; filename="{filename}"'})


# ── RACI Matrix Generator ─────────────────────────────────────────────────────

RACI_PROMPT = """You are a senior business analyst and process governance expert. Analyse the uploaded process documentation and produce a professional RACI Matrix document.

FORMAT RULES — follow exactly:
- Use markdown with ## for section headings
- Do NOT use strikethrough text (~~text~~) under any circumstances
- Do NOT use ALL CAPS text anywhere in the document. Use normal sentence case or title case only.
- Do NOT include a title or preamble — start directly with the first section
- ALL sections must be formatted as markdown tables — no prose paragraphs, no bullet lists
- Do NOT output any unstructured text blocks anywhere in the document

RACI MATRIX FORMAT:
| Activity | Role 1 | Role 2 | Role 3 | ... |
|---|---|---|---|---|
| **Section Name** | | | | |
| Activity name | R | A | C | ... |

Use only: R (Responsible), A (Accountable), C (Consulted), I (Informed). Each row must have exactly one A.
Group activities under bold section-header rows (all non-Activity cells empty).

ROLE GLOSSARY FORMAT — render as a markdown table:
| Role | Description | Typical Job Titles |
|---|---|---|
| Role name | One-sentence description of responsibilities | Title 1, Title 2 |

KEY FINDINGS FORMAT — render as a markdown table:
| Finding | Affected Role(s) | Recommended Action |
|---|---|---|
| Issue description | Role name | Action to take |

RECOMMENDATIONS FORMAT — render as a markdown table:
| # | Recommendation | Priority | Rationale |
|---|---|---|---|
| 1 | Action to take | High / Medium / Low | Why this matters |

SOURCES AND ASSUMPTIONS FORMAT — render as a markdown table:
| Ref | Assumption / Source | Basis |
|---|---|---|
| A1 | Statement of assumption | Where it came from |

ACRONYM DICTIONARY FORMAT — render as a markdown table:
| Term / Acronym | Definition |
|---|---|
| TERM | Plain-English definition |

STRUCTURE (include only selected sections):
- ## RACI Matrix
- ## Role Glossary
- ## Key Findings & Accountability Gaps
- ## Recommendations
- ## Sources and Assumptions
- ## Acronym & Term Dictionary — only if "glossary" is in the sections list

GRANULARITY GUIDANCE:
- high: 5–12 major activities (phase-level)
- detailed: 15–30 activities (function-level)
- task: 30–60 activities (individual task-level)

Be specific — use role names and activity names drawn directly from the uploaded documents."""


def _generate_raci(client: anthropic.Anthropic, document_text: str, process_name: str,
                   granularity: str, sections: list) -> str:
    sections_note = ", ".join(sections)
    user_msg = (
        f"Process name: {process_name}\n"
        f"Activity granularity: {granularity}\n"
        f"Sections to include: {sections_note}\n\n"
        f"Source documents:\n\n{document_text[:12000]}"
    )
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        system=RACI_PROMPT,
        messages=[{"role": "user", "content": user_msg}],
    )
    return msg.content[0].text


@app.post("/api/generate-raci")
async def generate_raci_endpoint(
    files:         List[UploadFile] = File(...),
    process_name:  str = Form(...),
    granularity:   str = Form("detailed"),
    sections_json: str = Form("[]"),
):
    if not files:
        raise HTTPException(status_code=400, detail="At least one file is required.")
    all_texts = []
    for f in files:
        ext = Path(f.filename).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise HTTPException(status_code=400, detail=f"Unsupported file type '{ext}'.")
        fb = await f.read()
        if len(fb) / 1024 / 1024 > MAX_FILE_SIZE_MB:
            raise HTTPException(status_code=413, detail=f"File '{f.filename}' too large.")
        try:
            t = extract_text(fb, f.filename)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Could not parse '{f.filename}': {e}")
        if t.strip():
            all_texts.append(f"=== {f.filename} ===\n{t.strip()}")
    if not all_texts:
        raise HTTPException(status_code=422, detail="No text could be extracted.")
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured.")
    sections = json.loads(sections_json) if sections_json.strip() else ['raci_matrix', 'role_glossary', 'sources']
    if 'raci_matrix' not in sections:
        sections = ['raci_matrix'] + sections
    if 'sources' not in sections:
        sections.append('sources')
    try:
        client = anthropic.Anthropic(api_key=api_key)
        md = _generate_raci(client, "\n\n".join(all_texts), process_name, granularity, sections)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"RACI generation failed: {e}")
    return JSONResponse({"status": "success", "raci_markdown": md})


@app.post("/api/export-raci")
async def export_raci(
    raci_markdown: str = Form(...),
    process_name:  str = Form("RACI Matrix"),
    format:        str = Form("docx"),
):
    safe_name = re.sub(r'[^a-z0-9_\-]', '_', process_name.lower()) or 'raci_matrix'
    fmt = format.lower().strip()
    try:
        if fmt == "pptx":
            content  = _make_bc_pptx(raci_markdown, process_name)
            media    = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename = f"{safe_name}.pptx"
        else:
            content  = _make_sop_docx(raci_markdown, process_name)
            media    = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            filename = f"{safe_name}.docx"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Export failed: {e}")
    return Response(content=content, media_type=media,
                    headers={"Content-Disposition": f'attachment; filename="{filename}"'})


# ── Change Impact Assessment ──────────────────────────────────────────────────

CIA_PROMPT = """You are a senior change management and process transformation consultant. Analyse the current state and future state process documents provided and produce a structured Change Impact Assessment.

FORMAT RULES — follow exactly:
- Use markdown with ## for section headings, ### for sub-headings, **bold** for labels
- Do NOT use strikethrough text (~~text~~) under any circumstances
- Do NOT use ALL CAPS text anywhere in the document. Use normal sentence case or title case only.
- Do NOT include a title or preamble — start directly with the first section
- All financial figures must be in US dollars ($)
- For every metric or assumption cited, add a superscript reference marker: value<sup>[A1]</sup>
- At the end, include a "## Sources and Assumptions" section listing every [A1], [A2]... marker with derivation, source, and caveats

STRUCTURE (include only selected sections):
- ## Executive Summary — 3–5 sentences: what is changing, overall impact magnitude, key recommendation
- ## Change Overview — what is driving the change, scope, timeline if known
- ## Stakeholder Impact by Group — for each stakeholder group:
  ### [Group Name] — Impact Level: [HIGH / MEDIUM / LOW]
  - What changes for them, skills/behaviour changes required, recommended engagement approach
- ## Process Delta (What Changes) — comparison of key process steps: what is removed, added, or modified
- ## People & Skills Impact — FTE changes, new skills required, training needs
- ## Technology & Systems Impact — systems added, retired, or changed; data migration considerations
- ## Risk & Change Readiness — top 3–5 change risks with likelihood, impact, and mitigation
- ## Recommended Change Actions — prioritised action plan: communication, training, transition management
- ## Sources and Assumptions — numbered [A1], [A2]... with derivation detail
- ## Acronym & Term Dictionary — if "glossary" is in the sections list, include an alphabetical dictionary of all acronyms, abbreviations, and industry-specific terms used in the document. Format each entry as **TERM / ACRONYM:** plain-English definition.

DEPTH GUIDANCE:
- summary: 2–3 bullets per section, high-level only
- standard: 4–6 bullets per section, balanced detail
- detailed: 8–12 bullets per section, comprehensive analysis

Base all findings on the actual content of the uploaded documents. Highlight gaps where information is limited."""


def _generate_cia(client: anthropic.Anthropic, current_text: str, future_text: str,
                  process_name: str, depth: str, sections: list) -> str:
    sections_note = ", ".join(sections)
    user_msg = (
        f"Process name: {process_name}\n"
        f"Analysis depth: {depth}\n"
        f"Sections to include: {sections_note}\n\n"
        f"=== CURRENT STATE DOCUMENTS ===\n{current_text[:6000]}\n\n"
        f"=== FUTURE STATE DOCUMENTS ===\n{future_text[:6000]}"
    )
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        system=CIA_PROMPT,
        messages=[{"role": "user", "content": user_msg}],
    )
    return msg.content[0].text


@app.post("/api/generate-cia")
async def generate_cia_endpoint(
    current_files: List[UploadFile] = File(...),
    future_files:  List[UploadFile] = File(...),
    process_name:  str = Form(...),
    depth:         str = Form("standard"),
    sections_json: str = Form("[]"),
):
    if not current_files:
        raise HTTPException(status_code=400, detail="At least one current state file is required.")
    if not future_files:
        raise HTTPException(status_code=400, detail="At least one future state file is required.")

    async def read_files(files):
        texts = []
        for f in files:
            ext = Path(f.filename).suffix.lower()
            if ext not in ALLOWED_EXTENSIONS:
                raise HTTPException(status_code=400, detail=f"Unsupported file type '{ext}'.")
            fb = await f.read()
            if len(fb) / 1024 / 1024 > MAX_FILE_SIZE_MB:
                raise HTTPException(status_code=413, detail=f"File '{f.filename}' too large.")
            try:
                t = extract_text(fb, f.filename)
            except Exception as e:
                raise HTTPException(status_code=422, detail=f"Could not parse '{f.filename}': {e}")
            if t.strip():
                texts.append(f"=== {f.filename} ===\n{t.strip()}")
        return texts

    current_texts = await read_files(current_files)
    future_texts  = await read_files(future_files)

    if not current_texts:
        raise HTTPException(status_code=422, detail="No text extracted from current state files.")
    if not future_texts:
        raise HTTPException(status_code=422, detail="No text extracted from future state files.")

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured.")

    sections = json.loads(sections_json) if sections_json.strip() else [
        'executive_summary', 'change_overview', 'stakeholder_impact', 'sources'
    ]
    if 'sources' not in sections:
        sections.append('sources')

    try:
        client = anthropic.Anthropic(api_key=api_key)
        md = _generate_cia(client, "\n\n".join(current_texts), "\n\n".join(future_texts),
                           process_name, depth, sections)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"CIA generation failed: {e}")
    return JSONResponse({"status": "success", "cia_markdown": md})


@app.post("/api/export-cia")
async def export_cia(
    cia_markdown:  str = Form(...),
    process_name:  str = Form("Change Impact Assessment"),
    format:        str = Form("docx"),
):
    safe_name = re.sub(r'[^a-z0-9_\-]', '_', process_name.lower()) or 'change_impact'
    fmt = format.lower().strip()
    try:
        if fmt == "pptx":
            content  = _make_bc_pptx(cia_markdown, process_name)
            media    = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename = f"{safe_name}.pptx"
        else:
            content  = _make_sop_docx(cia_markdown, process_name)
            media    = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            filename = f"{safe_name}.docx"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Export failed: {e}")
    return Response(content=content, media_type=media,
                    headers={"Content-Disposition": f'attachment; filename="{filename}"'})


# ── Serve frontend static files (must come LAST — after all API routes) ───────
from pathlib import Path as _Path
from fastapi.responses import FileResponse as _FileResponse
_PUBLIC = (_Path(__file__).resolve().parent.parent / "public").resolve()

@app.get("/{filename:path}")
async def serve_frontend(filename: str):
    """Serve any file from the public/ folder so the frontend works on localhost."""
    if not filename or filename == "/":
        filename = "index.html"
    target = (_PUBLIC / filename).resolve()
    # Security: ensure resolved path is inside _PUBLIC (use Path.is_relative_to for Windows safety)
    try:
        target.relative_to(_PUBLIC)
    except ValueError:
        raise HTTPException(status_code=403)
    if target.is_file():
        return _FileResponse(str(target))
    # Fall back to index.html for SPA-style navigation
    index = _PUBLIC / "index.html"
    if index.is_file():
        return _FileResponse(str(index))
    raise HTTPException(status_code=404)
